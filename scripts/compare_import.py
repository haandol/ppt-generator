"""원본 PPTX ↔ 임포트된 DesignSpec 차이 평가 도구.

원본 PPTX의 텍스트 요소를 1280×720 px 좌표로 정규화한 뒤, ImportService 가
생성한 DesignSpec 의 대응 요소와 매칭하여 위치/폰트크기/텍스트 차이를 슬라이드별로
리포트한다. 임포트 품질을 매 수정마다 정량적으로 추적하는 용도.

사용법:
    uv run python scripts/compare_import.py "<원본.pptx>" [--slides 2,19] [--json]
"""

from __future__ import annotations

import argparse
import json
import logging
import sys
from dataclasses import dataclass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from ppt_generator.interfaces.constants import (
    IMPORT_EMU_TO_PX,
    SLIDES_HEIGHT_PX,
    SLIDES_WIDTH_PX,
)
from ppt_generator.tools.pptx_import.service import ImportService

logging.disable(logging.CRITICAL)

# 매칭/차이 임계값 (px, pt)
POS_TOLERANCE_PX = 8.0
SIZE_TOLERANCE_PT = 2.0


@dataclass
class TextElement:
    """텍스트 요소 하나 (원본 또는 임포트 공통 표현)."""

    text: str
    left: float
    top: float
    width: float
    height: float
    font_size: float | None

    @property
    def norm_text(self) -> str:
        return "".join(self.text.split()).lower()


def _scale(prs: Presentation) -> tuple[float, float]:
    sw = prs.slide_width * IMPORT_EMU_TO_PX
    sh = prs.slide_height * IMPORT_EMU_TO_PX
    return (
        SLIDES_WIDTH_PX / sw if sw else 1.0,
        SLIDES_HEIGHT_PX / sh if sh else 1.0,
    )


def _max_font_pt(shape) -> float | None:
    """shape 텍스트의 최대 명시 폰트 크기(pt). 없으면 None."""
    best: float | None = None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rpr = run._r.find(qn("a:rPr"))
            sz = rpr.get("sz") if rpr is not None else None
            if sz:
                pt = int(sz) / 100
                if best is None or pt > best:
                    best = pt
    return best


def extract_original(prs: Presentation, slide_idx: int) -> list[TextElement]:
    """원본 슬라이드의 텍스트 요소를 px 좌표로 추출 (그룹 변환 반영)."""
    sx, sy = _scale(prs)
    out: list[TextElement] = []

    def walk(shapes, ox=0.0, oy=0.0, gsx=1.0, gsy=1.0, cox=0.0, coy=0.0):
        for s in shapes:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                gx = s._element.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
                nox, noy, ngsx, ngsy, ncox, ncoy = ox, oy, gsx, gsy, cox, coy
                if gx is not None:
                    off = gx.find(qn("a:off"))
                    ext = gx.find(qn("a:ext"))
                    ch_off = gx.find(qn("a:chOff"))
                    ch_ext = gx.find(qn("a:chExt"))
                    if None not in (off, ext, ch_off, ch_ext):
                        nox = int(off.get("x", "0")) * IMPORT_EMU_TO_PX * sx
                        noy = int(off.get("y", "0")) * IMPORT_EMU_TO_PX * sy
                        ecx, ecy = int(ext.get("cx", "1")), int(ext.get("cy", "1"))
                        ccx, ccy = (
                            int(ch_ext.get("cx", "1")),
                            int(ch_ext.get("cy", "1")),
                        )
                        ngsx = ecx / ccx if ccx else 1.0
                        ngsy = ecy / ccy if ccy else 1.0
                        ncox = int(ch_off.get("x", "0")) * IMPORT_EMU_TO_PX * sx
                        ncoy = int(ch_off.get("y", "0")) * IMPORT_EMU_TO_PX * sy
                walk(s.shapes, nox, noy, ngsx, ngsy, ncox, ncoy)
                continue
            if not getattr(s, "has_text_frame", False):
                continue
            if not s.text_frame.text.strip():
                continue
            left = s.left * IMPORT_EMU_TO_PX * sx
            top = s.top * IMPORT_EMU_TO_PX * sy
            w = s.width * IMPORT_EMU_TO_PX * sx
            h = s.height * IMPORT_EMU_TO_PX * sy
            # 그룹 내부면 자식 좌표계 → 부모 좌표계 변환
            if gsx != 1.0 or gsy != 1.0 or cox or coy or ox or oy:
                left = (left - cox) * gsx + ox
                top = (top - coy) * gsy + oy
                w *= gsx
                h *= gsy
            out.append(
                TextElement(
                    text=s.text_frame.text,
                    left=round(left, 1),
                    top=round(top, 1),
                    width=round(w, 1),
                    height=round(h, 1),
                    font_size=_max_font_pt(s),
                )
            )

    walk(prs.slides[slide_idx].shapes)

    # 레이아웃/마스터의 정적 텍스트(로고 캡션, "Thank you!", 저작권 등)도 포함한다.
    # PowerPoint 는 이들을 슬라이드에 렌더하며, import 도 이를 상속하기 때문이다.
    seen = {e.norm_text for e in out}
    slide = prs.slides[slide_idx]
    for src in (
        getattr(slide, "slide_layout", None),
        getattr(getattr(slide, "slide_layout", None), "slide_master", None),
    ):
        if src is None:
            continue
        for s in src.shapes:
            if getattr(s, "is_placeholder", False):
                continue
            if not getattr(s, "has_text_frame", False):
                continue
            if not s.text_frame.text.strip():
                continue
            e = TextElement(
                text=s.text_frame.text,
                left=round(s.left * IMPORT_EMU_TO_PX * sx, 1),
                top=round(s.top * IMPORT_EMU_TO_PX * sy, 1),
                width=round(s.width * IMPORT_EMU_TO_PX * sx, 1),
                height=round(s.height * IMPORT_EMU_TO_PX * sy, 1),
                font_size=_max_font_pt(s),
            )
            if e.norm_text in seen:
                continue
            seen.add(e.norm_text)
            out.append(e)
    return out


def extract_imported(spec_slide) -> list[TextElement]:
    """임포트된 슬라이드 spec 의 텍스트박스/도형 텍스트 요소 추출."""
    out: list[TextElement] = []
    for tb in spec_slide.textboxes:
        text = "\n".join("".join(r.text for r in p.runs) for p in tb.paragraphs)
        if not text.strip():
            continue
        fs = None
        for p in tb.paragraphs:
            for r in p.runs:
                if r.font_size_pt and (fs is None or r.font_size_pt > fs):
                    fs = r.font_size_pt
        out.append(
            TextElement(text, tb.left_px, tb.top_px, tb.width_px, tb.height_px, fs)
        )
    for sh in spec_slide.shapes:
        text = sh.text or "\n".join(
            "".join(r.text for r in p.runs) for p in sh.paragraphs
        )
        if not text or not text.strip():
            continue
        fs = sh.text_size_pt
        for p in sh.paragraphs:
            for r in p.runs:
                if r.font_size_pt and (fs is None or r.font_size_pt > fs):
                    fs = r.font_size_pt
        out.append(
            TextElement(text, sh.left_px, sh.top_px, sh.width_px, sh.height_px, fs)
        )
    return out


def match_and_diff(
    originals: list[TextElement], imported: list[TextElement]
) -> tuple[list[dict], list[TextElement], list[TextElement]]:
    """정규화 텍스트로 매칭 후 위치/폰트 차이 계산."""
    remaining = list(imported)
    diffs: list[dict] = []
    unmatched_orig: list[TextElement] = []
    for o in originals:
        cand = [i for i in remaining if i.norm_text == o.norm_text]
        if not cand:
            # 부분 포함 매칭 (줄바꿈/공백 차이 흡수)
            cand = [
                i
                for i in remaining
                if o.norm_text
                and (o.norm_text in i.norm_text or i.norm_text in o.norm_text)
            ]
        if not cand:
            unmatched_orig.append(o)
            continue
        m = min(cand, key=lambda i: abs(i.left - o.left) + abs(i.top - o.top))
        remaining.remove(m)
        dl, dt = m.left - o.left, m.top - o.top
        dfs = (
            (m.font_size - o.font_size)
            if (m.font_size is not None and o.font_size is not None)
            else None
        )
        diffs.append(
            {
                "text": o.text[:40].replace("\n", "⏎"),
                "orig_pos": (o.left, o.top),
                "dpos": (round(dl, 1), round(dt, 1)),
                "orig_font": o.font_size,
                "imp_font": m.font_size,
                "dfont": round(dfs, 1) if dfs is not None else None,
                "pos_ok": abs(dl) <= POS_TOLERANCE_PX and abs(dt) <= POS_TOLERANCE_PX,
                "font_ok": dfs is None or abs(dfs) <= SIZE_TOLERANCE_PT,
            }
        )
    return diffs, unmatched_orig, remaining


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--slides", default="", help="1-based, e.g. 2,19")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    spec, _ = ImportService().import_from_file(args.pptx)
    total = len(prs.slides)
    which = (
        [int(x) - 1 for x in args.slides.split(",") if x.strip()]
        if args.slides
        else list(range(total))
    )

    report = []
    grand = {
        "matched": 0,
        "pos_fail": 0,
        "font_fail": 0,
        "unmatched_o": 0,
        "extra_i": 0,
    }
    for idx in which:
        o = extract_original(prs, idx)
        i = extract_imported(spec.slides[idx])
        diffs, unmatched_o, extra_i = match_and_diff(o, i)
        pos_fail = sum(1 for d in diffs if not d["pos_ok"])
        font_fail = sum(1 for d in diffs if not d["font_ok"])
        grand["matched"] += len(diffs)
        grand["pos_fail"] += pos_fail
        grand["font_fail"] += font_fail
        grand["unmatched_o"] += len(unmatched_o)
        grand["extra_i"] += len(extra_i)
        report.append(
            {
                "slide": idx + 1,
                "matched": len(diffs),
                "pos_fail": pos_fail,
                "font_fail": font_fail,
                "unmatched_original": [
                    u.text[:40].replace("\n", "⏎") for u in unmatched_o
                ],
                "extra_imported": [e.text[:40].replace("\n", "⏎") for e in extra_i],
                "diffs": diffs,
            }
        )

    if args.json:
        print(
            json.dumps(
                {"summary": grand, "slides": report}, ensure_ascii=False, indent=2
            )
        )
        return 0

    for r in report:
        flag = "✅" if (r["pos_fail"] == 0 and r["font_fail"] == 0) else "⚠️"
        print(
            f"\n{flag} SLIDE {r['slide']}: matched={r['matched']} "
            f"pos_fail={r['pos_fail']} font_fail={r['font_fail']} "
            f"unmatched_orig={len(r['unmatched_original'])} extra_imp={len(r['extra_imported'])}"
        )
        for d in r["diffs"]:
            if d["pos_ok"] and d["font_ok"]:
                continue
            marks = []
            if not d["pos_ok"]:
                marks.append(f"Δpos={d['dpos']}")
            if not d["font_ok"]:
                marks.append(f"font {d['orig_font']}→{d['imp_font']} (Δ{d['dfont']})")
            print(f"    ⚠️ {d['text']!r}: {', '.join(marks)}")
        for u in r["unmatched_original"]:
            print(f"    ❌ 원본에만 있음(누락): {u!r}")

    g = grand
    print(
        f"\n{'=' * 60}\nTOTAL: matched={g['matched']} "
        f"pos_fail={g['pos_fail']} font_fail={g['font_fail']} "
        f"원본누락={g['unmatched_o']} 임포트초과={g['extra_i']}"
    )
    score = (
        100 * (g["matched"] - g["pos_fail"] - g["font_fail"]) / g["matched"]
        if g["matched"]
        else 0
    )
    print(f"위치+폰트 정확도: {score:.1f}% (매칭 요소 기준)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
