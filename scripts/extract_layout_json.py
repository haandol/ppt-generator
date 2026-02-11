"""PPTX 템플릿에서 전체 레이아웃 정보를 layout.json으로 추출.

사용법:
    uv run python scripts/extract_layout_json.py

출력:
    template/layout.json — 모든 레이아웃의 placeholder/shape 위치·서식 정보
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 템플릿 경로 (프로젝트 루트 기준)
PROJECT_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE_PATH = PROJECT_ROOT / "template" / "template.pptx"
OUTPUT_PATH = PROJECT_ROOT / "template" / "layout.json"

# 슬라이드 크기 (EMU → px 변환)
SLIDE_WIDTH_EMU = 12_192_000
SLIDE_HEIGHT_EMU = 6_858_000
HTML_WIDTH_PX = 1280
HTML_HEIGHT_PX = 720

# 메타데이터 placeholder 타입 (python-pptx enum int 값)
# DATE=16, FOOTER=15, SLIDE_NUMBER=13
METADATA_PH_TYPES = {13, 15, 16}


def emu_to_px_x(emu: int) -> int:
    return round(emu * HTML_WIDTH_PX / SLIDE_WIDTH_EMU)


def emu_to_px_y(emu: int) -> int:
    return round(emu * HTML_HEIGHT_PX / SLIDE_HEIGHT_EMU)


def emu_to_px_rect(left: int, top: int, width: int, height: int) -> dict:
    return {
        "left": emu_to_px_x(left),
        "top": emu_to_px_y(top),
        "width": emu_to_px_x(width),
        "height": emu_to_px_y(height),
    }


def _color_to_hex(color) -> str | None:
    """python-pptx color 객체에서 hex 문자열 추출."""
    try:
        if color is None:
            return None
        if hasattr(color, "rgb") and color.rgb is not None:
            return f"#{color.rgb}"
        if hasattr(color, "theme_color") and color.theme_color is not None:
            return f"theme:{color.theme_color}"
    except Exception:
        pass
    return None


def _fill_to_dict(fill) -> dict | None:
    """python-pptx FillFormat에서 색상 정보 추출."""
    try:
        if fill is None:
            return None
        fill_type = fill.type
        if fill_type is None:
            return None
        type_name = str(fill_type)
        if "SOLID" in type_name:
            fg = _color_to_hex(fill.fore_color)
            return {"type": "solid", "color": fg}
        if "GRADIENT" in type_name:
            return {"type": "gradient"}
        if "PATTERN" in type_name:
            return {"type": "pattern"}
        if "PICTURE" in type_name:
            return {"type": "picture"}
    except Exception:
        pass
    return None


def _extract_placeholder(ph) -> dict:
    """placeholder에서 정보 추출."""
    ph_type_int = int(ph.placeholder_format.type)
    ph_type_name = str(ph.placeholder_format.type).split("(")[0].strip()
    if " " in ph_type_name:
        ph_type_name = ph_type_name.split(" ")[-1]

    is_metadata = ph_type_int in METADATA_PH_TYPES

    result = {
        "ph_idx": ph.placeholder_format.idx,
        "ph_type": ph_type_int,
        "ph_type_name": ph_type_name,
        "name": ph.name,
        "is_metadata": is_metadata,
        "position": emu_to_px_rect(ph.left, ph.top, ph.width, ph.height),
    }

    return result


def _extract_shape(shape) -> dict | None:
    """비-placeholder shape에서 정보 추출."""
    # placeholder는 별도 처리
    if shape.is_placeholder:
        return None

    shape_type = shape.shape_type
    type_name = str(shape_type) if shape_type else "unknown"

    result = {
        "name": shape.name,
        "shape_type": type_name,
        "position": emu_to_px_rect(shape.left, shape.top, shape.width, shape.height),
    }

    # 도형 fill
    try:
        if hasattr(shape, "fill"):
            fill_info = _fill_to_dict(shape.fill)
            if fill_info:
                result["fill"] = fill_info
    except Exception:
        pass

    # 텍스트 포함 여부
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            result["text"] = text

    return result


def _detect_theme(layout_name: str, layout_index: int) -> str:
    """레이아웃 이름과 인덱스로 테마(light/dark) 판별."""
    name_lower = layout_name.lower()
    if "gradient" in name_lower or "dark" in name_lower:
        return "dark"
    # 인덱스 96은 Title Only Dark
    if layout_index == 96:
        return "dark"
    # 인덱스 28은 Gradient 포함 레이아웃
    if layout_index == 28:
        return "dark"
    return "light"


def _classify_role(ph_type_int: int, ph_idx: int, ph_name: str) -> str | None:
    """placeholder의 역할(title/subtitle/body/picture)을 분류."""
    # PP_PLACEHOLDER 타입 enum:
    # TITLE=1, CENTER_TITLE=3 → title
    # SUBTITLE=4 → subtitle
    # BODY=2, OBJECT=7 → body
    # PICTURE=18 → picture
    if ph_type_int in (1, 3):
        return "title"
    if ph_type_int == 4:
        return "subtitle"
    if ph_type_int in (2, 7):
        return "body"
    if ph_type_int == 18:
        return "picture"
    return None


def extract_all_layouts(template_path: Path) -> dict:
    """전체 레이아웃 정보 추출."""
    prs = Presentation(str(template_path))
    layouts: list[dict] = []

    for i, layout in enumerate(prs.slide_layouts):
        name = layout.name

        # "Do Not Use" 레이아웃 표시
        is_excluded = "do not use" in name.lower()

        theme = _detect_theme(name, i)

        placeholders = []
        for ph in layout.placeholders:
            ph_data = _extract_placeholder(ph)
            # 역할 분류
            role = _classify_role(ph_data["ph_type"], ph_data["ph_idx"], ph_data["name"])
            if role:
                ph_data["role"] = role
            placeholders.append(ph_data)

        shapes = []
        for shape in layout.shapes:
            shape_data = _extract_shape(shape)
            if shape_data:
                shapes.append(shape_data)

        # 주요 영역 좌표 요약 (HTML 스켈레톤에서 사용)
        regions = {}
        for ph_data in placeholders:
            if ph_data.get("is_metadata"):
                continue
            role = ph_data.get("role")
            if role and role not in regions:
                regions[role] = ph_data["position"]

        layout_data: dict = {
            "layout_index": i,
            "layout_name": name,
            "theme": theme,
        }
        if is_excluded:
            layout_data["excluded"] = True

        layout_data["regions"] = regions
        layout_data["placeholders"] = placeholders

        if shapes:
            layout_data["shapes"] = shapes

        layouts.append(layout_data)

    return {
        "slide_size": {"width": HTML_WIDTH_PX, "height": HTML_HEIGHT_PX},
        "total_layouts": len(layouts),
        "layouts": layouts,
    }


def main() -> None:
    if not TEMPLATE_PATH.exists():
        print(f"Error: Template not found at {TEMPLATE_PATH}")
        return

    print(f"Template: {TEMPLATE_PATH}")
    data = extract_all_layouts(TEMPLATE_PATH)
    print(f"Extracted {data['total_layouts']} layouts")

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    with open(OUTPUT_PATH, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

    print(f"Output: {OUTPUT_PATH}")
    print(f"File size: {OUTPUT_PATH.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
