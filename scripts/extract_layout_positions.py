"""PPTX 템플릿에서 레이아웃별 placeholder 위치를 추출하여 1280x720px HTML 좌표로 변환.

사용법:
    uv run python scripts/extract_layout_positions.py

출력:
    각 레이아웃별 placeholder 위치를 JSON 형식으로 콘솔 출력
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# AWS PPTX 템플릿 경로
TEMPLATE_PATH = Path.home() / "Documents" / "work" / "ppt-generator" / "2026 Confidential AWS Powerpoint Template Light & Dark Themes.pptx"

# LAYOUT_MAP에서 가져온 레이아웃 인덱스
LAYOUT_MAP: dict[str, int] = {
    "title": 0,
    "text_image": 28,
    "text_only": 22,
    "chart": 21,
    "closing": 87,
    "freeform": 88,
}

# 슬라이드 물리 크기 (인치)
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5

# HTML 슬라이드 크기 (px)
HTML_WIDTH_PX = 1280
HTML_HEIGHT_PX = 720

# 변환 계수
INCHES_TO_PX_X = HTML_WIDTH_PX / SLIDE_WIDTH_INCHES  # ~96.0
INCHES_TO_PX_Y = HTML_HEIGHT_PX / SLIDE_HEIGHT_INCHES  # ~96.0


def emu_to_inches(emu: int) -> float:
    """EMU 값을 인치로 변환."""
    return emu / 914400


def inches_to_px(inches_x: float, inches_y: float) -> tuple[int, int]:
    """인치 좌표를 px로 변환."""
    return (
        round(inches_x * INCHES_TO_PX_X),
        round(inches_y * INCHES_TO_PX_Y),
    )


def extract_layout_regions(template_path: Path) -> dict:
    """템플릿에서 각 레이아웃의 placeholder 위치를 추출."""
    prs = Presentation(str(template_path))
    results: dict[str, dict] = {}

    for layout_type, layout_index in LAYOUT_MAP.items():
        layout = prs.slide_layouts[layout_index]
        layout_data = {
            "layout_name": layout.name,
            "layout_index": layout_index,
            "placeholders": [],
        }

        for ph in layout.placeholders:
            left_inches = emu_to_inches(ph.left)
            top_inches = emu_to_inches(ph.top)
            width_inches = emu_to_inches(ph.width)
            height_inches = emu_to_inches(ph.height)

            left_px, top_px = inches_to_px(left_inches, top_inches)
            width_px = round(width_inches * INCHES_TO_PX_X)
            height_px = round(height_inches * INCHES_TO_PX_Y)

            layout_data["placeholders"].append({
                "ph_idx": ph.placeholder_format.idx,
                "ph_type": str(ph.placeholder_format.type),
                "name": ph.name,
                "position_inches": {
                    "left": round(left_inches, 3),
                    "top": round(top_inches, 3),
                    "width": round(width_inches, 3),
                    "height": round(height_inches, 3),
                },
                "position_px": {
                    "left": left_px,
                    "top": top_px,
                    "width": width_px,
                    "height": height_px,
                },
            })

        results[layout_type] = layout_data

    return results


def format_layout_regions_for_constants(results: dict) -> dict[str, dict]:
    """추출 결과를 LAYOUT_REGIONS 상수 형태로 정리.

    각 레이아웃에서 의미 있는 영역(title, subtitle, body, image)을
    placeholder 타입 기반으로 매핑.
    """
    regions: dict[str, dict] = {}

    for layout_type, data in results.items():
        region: dict[str, dict[str, int]] = {}
        for ph in data["placeholders"]:
            px = ph["position_px"]
            ph_type = ph["ph_type"]
            ph_name = ph["name"].lower()

            area = {"left": px["left"], "top": px["top"], "width": px["width"], "height": px["height"]}

            if "CENTER_TITLE" in ph_type or "TITLE" in ph_type:
                if "subtitle" in ph_name or "SUBTITLE" in ph_type:
                    region["subtitle"] = area
                else:
                    region["title"] = area
            elif "SUBTITLE" in ph_type:
                region["subtitle"] = area
            elif "BODY" in ph_type or "OBJECT" in ph_type:
                region["body"] = area
            elif "PICTURE" in ph_type:
                region["image"] = area

        regions[layout_type] = region

    return regions


def main() -> None:
    if not TEMPLATE_PATH.exists():
        print(f"Error: Template not found at {TEMPLATE_PATH}")
        return

    print(f"Template: {TEMPLATE_PATH}\n")

    # 상세 추출 결과
    results = extract_layout_regions(TEMPLATE_PATH)
    print("=== Detailed Placeholder Positions ===")
    print(json.dumps(results, indent=2, ensure_ascii=False))

    # LAYOUT_REGIONS 상수용 정리
    regions = format_layout_regions_for_constants(results)
    print("\n=== LAYOUT_REGIONS (for constants.py) ===")
    print(json.dumps(regions, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
