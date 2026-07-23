"""HTML/PPTX 렌더 패리티 계약 테스트."""

from __future__ import annotations

from dataclasses import fields, is_dataclass
from typing import get_args, get_origin, get_type_hints

import pytest
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

from ppt_generator.interfaces.constants import (
    EMU_PER_INCH,
    EXPORT_PX_TO_INCHES_X,
    EXPORT_PX_TO_INCHES_Y,
    PX_TO_EMU,
)
from ppt_generator.interfaces.render_parity import (
    RENDER_FIELD_POLICIES,
    RenderFieldKind,
)
from ppt_generator.interfaces.schemas import (
    DesignSpec,
    PptxParagraph,
    PptxShape,
    PptxSlideSpec,
    PptxTextBox,
    PptxTextRun,
)
from ppt_generator.tools.pptx.service import ExportService
from ppt_generator.tools.slides.html_renderer import spec_to_html_section


def _nested_dataclasses(annotation: object) -> set[type]:
    origin = get_origin(annotation)
    if origin is not None:
        nested: set[type] = set()
        for argument in get_args(annotation):
            nested.update(_nested_dataclasses(argument))
        return nested
    if isinstance(annotation, type) and is_dataclass(annotation):
        return {annotation}
    return set()


def _reachable_schema_types(root: type) -> set[type]:
    pending = [root]
    discovered: set[type] = set()
    while pending:
        schema_type = pending.pop()
        if schema_type in discovered:
            continue
        discovered.add(schema_type)
        for annotation in get_type_hints(schema_type).values():
            pending.extend(_nested_dataclasses(annotation) - discovered)
    return discovered


def _px_x(value: int) -> float:
    return value / EMU_PER_INCH / EXPORT_PX_TO_INCHES_X


def _px_y(value: int) -> float:
    return value / EMU_PER_INCH / EXPORT_PX_TO_INCHES_Y


def _parity_spec() -> DesignSpec:
    return DesignSpec(
        slides=[
            PptxSlideSpec(
                background_color="#123456",
                speaker_notes="Parity notes",
                shapes=[
                    PptxShape(
                        left_px=40,
                        top_px=50,
                        width_px=280,
                        height_px=120,
                        shape_type="rounded_rectangle",
                        fill_color="#DDEEFF",
                        border_color="#234567",
                        border_width_pt=2,
                        corner_radius_px=10,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="Card body",
                                        font_size_pt=16,
                                        color="#112233",
                                        bold=True,
                                    )
                                ],
                                alignment="center",
                            )
                        ],
                        vertical_alignment="middle",
                        padding_left_px=12,
                        padding_right_px=10,
                        padding_top_px=8,
                        padding_bottom_px=6,
                        rotation=15,
                        z_index=0,
                    ),
                    PptxShape(
                        left_px=360,
                        top_px=110,
                        width_px=180,
                        height_px=60,
                        shape_type="line",
                        border_color="#CC5500",
                        border_width_pt=2,
                        dash_style="dash",
                        start_arrow=True,
                        end_arrow=True,
                        z_index=1,
                    ),
                ],
                textboxes=[
                    PptxTextBox(
                        left_px=80,
                        top_px=220,
                        width_px=520,
                        height_px=180,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="Primary",
                                        font_size_pt=14,
                                        color="#445566",
                                        bold=True,
                                        italic=True,
                                        font_name="Arial",
                                        href="https://example.com",
                                    )
                                ],
                                alignment="left",
                                space_before_pt=3,
                                space_after_pt=5,
                            ),
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="Secondary",
                                        font_size_pt=14,
                                        color="#445566",
                                    )
                                ],
                                bullet_level=0,
                                alignment="left",
                            ),
                        ],
                        line_spacing_pt=24,
                        vertical_alignment="middle",
                        padding_left_px=9,
                        padding_right_px=11,
                        padding_top_px=7,
                        padding_bottom_px=5,
                        autofit="none",
                        z_index=2,
                    )
                ],
            )
        ]
    )


class TestRenderFieldClassification:
    def test_every_design_spec_field_has_exactly_one_policy(self):
        reachable = _reachable_schema_types(DesignSpec)

        assert set(RENDER_FIELD_POLICIES) == reachable
        for schema_type in reachable:
            expected = {field.name for field in fields(schema_type)}
            actual = set(RENDER_FIELD_POLICIES[schema_type])
            assert actual == expected, schema_type.__name__

    def test_every_non_shared_field_has_a_rationale(self):
        for schema_type, policies in RENDER_FIELD_POLICIES.items():
            for field_name, policy in policies.items():
                if policy.kind is RenderFieldKind.SHARED:
                    continue
                assert policy.rationale.strip(), (
                    f"{schema_type.__name__}.{field_name} needs a rationale"
                )


class TestCrossRenderContract:
    def test_shared_spec_semantics_are_present_in_html_and_pptx(self, tmp_path):
        design_spec = _parity_spec()
        slide_spec = design_spec.slides[0]

        html = spec_to_html_section(0, slide_spec)
        response = ExportService().export_from_design_spec(
            design_spec,
            output_dir=tmp_path,
            bg_image_policy="none",
        )
        slide = Presentation(response.pptx_path).slides[0]

        assert "background-color:#123456" in html
        assert 'data-speaker-notes="Parity notes"' in html
        assert "left:40px;top:50px;width:280px;height:120px" in html
        assert "background-color:#DDEEFF" in html
        assert "border:2pt solid #234567" in html
        assert "border-radius:10px" in html
        assert "transform:rotate(15deg)" in html
        assert "padding:8px 10px 6px 12px" in html
        assert "justify-content:center" in html
        assert 'stroke="#CC5500"' in html
        assert 'stroke-dasharray="8,6"' in html
        assert 'marker-start="url(#ah-start)"' in html
        assert 'marker-end="url(#ah-end)"' in html
        assert "left:80px;top:220px;width:520px;height:180px" in html
        assert "padding:7px 11px 5px 9px" in html
        assert html.count("line-height:24pt") >= 3
        assert "padding-top:3pt" in html
        assert "padding-bottom:5pt" in html
        assert "font-size:14.00pt" in html
        assert "color:#445566" in html
        assert "font-weight:bold" in html
        assert "font-style:italic" in html
        assert "font-family:'Arial'" in html
        assert 'href="https://example.com"' in html
        assert "<li " in html
        assert "box-shadow" not in html

        assert str(slide.background.fill.fore_color.rgb) == "123456"
        assert slide.notes_slide.notes_text_frame.text == "Parity notes"

        card = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and shape.text == "Card body"
        )
        textbox = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Primary" in shape.text
        )
        connector = next(
            element for element in slide.shapes._spTree if element.tag == qn("p:cxnSp")
        )

        assert _px_x(card.left) == pytest.approx(40, abs=0.01)
        assert _px_y(card.top) == pytest.approx(50, abs=0.01)
        assert _px_x(card.width) == pytest.approx(280, abs=0.01)
        assert _px_y(card.height) == pytest.approx(120, abs=0.01)
        assert str(card.fill.fore_color.rgb) == "DDEEFF"
        assert str(card.line.color.rgb) == "234567"
        assert card.line.width.pt == pytest.approx(2)
        assert card.rotation == pytest.approx(15)
        assert card.text_frame.margin_left == pytest.approx(12 * PX_TO_EMU)
        assert card.text_frame.margin_right == pytest.approx(10 * PX_TO_EMU)
        assert card.text_frame.margin_top == pytest.approx(8 * PX_TO_EMU)
        assert card.text_frame.margin_bottom == pytest.approx(6 * PX_TO_EMU)
        assert card.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
        assert card.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER

        line = connector.find(qn("p:spPr")).find(qn("a:ln"))
        assert line.find(qn("a:prstDash")).get("val") == "dash"
        assert line.find(qn("a:headEnd")).get("type") == "triangle"
        assert line.find(qn("a:tailEnd")).get("type") == "triangle"

        assert _px_x(textbox.left) == pytest.approx(80, abs=0.01)
        assert _px_y(textbox.top) == pytest.approx(220, abs=0.01)
        assert _px_x(textbox.width) == pytest.approx(520, abs=0.01)
        assert _px_y(textbox.height) == pytest.approx(180, abs=0.01)
        assert textbox.text_frame.margin_left == pytest.approx(9 * PX_TO_EMU)
        assert textbox.text_frame.margin_right == pytest.approx(11 * PX_TO_EMU)
        assert textbox.text_frame.margin_top == pytest.approx(7 * PX_TO_EMU)
        assert textbox.text_frame.margin_bottom == pytest.approx(5 * PX_TO_EMU)
        assert textbox.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE

        primary = textbox.text_frame.paragraphs[0]
        assert primary.alignment == PP_ALIGN.LEFT
        assert primary.space_before.pt == pytest.approx(3)
        assert primary.space_after.pt == pytest.approx(5)
        assert primary.line_spacing.pt == pytest.approx(24)
        run = primary.runs[0]
        assert run.font.size.pt == pytest.approx(14)
        assert run.font.bold is True
        assert run.font.italic is True
        assert run.font.name == "Arial"
        assert str(run.font.color.rgb) == "445566"
        assert run.hyperlink.address == "https://example.com"

        bullet = textbox.text_frame.paragraphs[1]
        assert bullet._p.find(f".//{qn('a:buChar')}") is not None

        effect_refs = [
            effect_ref.get("idx")
            for element in (card._element, connector)
            if (style := element.find(qn("p:style"))) is not None
            if (effect_ref := style.find(qn("a:effectRef"))) is not None
        ]
        assert effect_refs == ["0", "0"]
