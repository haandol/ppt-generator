"""Pytest 공통 fixture.

데이터 생성 helper 는 `_helpers.py` 에 두고 각 테스트 파일에서
`from _helpers import ...` 로 직접 import 한다 (`tests/` 가 pytest
rootdir 에 의해 sys.path 에 들어와 있어 가능).
"""

from __future__ import annotations

import pytest
from pptx import Presentation


@pytest.fixture
def blank_slide():
    """python-pptx 빈 슬라이드 fixture."""
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])
