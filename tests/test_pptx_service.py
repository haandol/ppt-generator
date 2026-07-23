from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from ppt_generator.interfaces.schemas import (
    DesignSpec,
    PptxParagraph,
    PptxShape,
    PptxSlideSpec,
    PptxTextBox,
    PptxTextRun,
)
from ppt_generator.tools.pptx.service import ExportService


def _make_design_spec() -> DesignSpec:
    return DesignSpec(
        slides=[
            PptxSlideSpec(
                background_color="#1a1a2e",
                textboxes=[
                    PptxTextBox(
                        left_px=40,
                        top_px=40,
                        width_px=600,
                        height_px=60,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="제목",
                                        font_size_pt=32,
                                        bold=True,
                                        color="#ffffff",
                                    )
                                ]
                            )
                        ],
                    ),
                    PptxTextBox(
                        left_px=40,
                        top_px=120,
                        width_px=600,
                        height_px=400,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="본문 텍스트",
                                        font_size_pt=18,
                                        color="#cccccc",
                                    )
                                ]
                            )
                        ],
                    ),
                ],
                shapes=[
                    PptxShape(
                        left_px=700,
                        top_px=120,
                        width_px=500,
                        height_px=400,
                        fill_color="#2a2a4e",
                        shape_type="rounded_rectangle",
                        text="도형 텍스트",
                        text_color="#ffffff",
                    ),
                ],
                speaker_notes="발표자 노트",
            ),
        ]
    )


@pytest.fixture
def service():
    return ExportService()


class TestExportFromDesignSpec:
    def test_creates_pptx_file(self, service, tmp_path):
        spec = _make_design_spec()
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        path = Path(response.pptx_path)
        assert path.exists()
        assert path.suffix == ".pptx"

    def test_raises_on_empty_slides(self, service):
        spec = DesignSpec(slides=[])
        with pytest.raises(ValueError, match="디자인 스펙에 슬라이드가 없습니다"):
            service.export_from_design_spec(spec)

    def test_extracts_text(self, service, tmp_path):
        spec = _make_design_spec()
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        all_text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
        assert "제목" in all_text
        assert "본문 텍스트" in all_text

    def test_preserves_paragraph_spacing_for_textboxes_and_shapes(
        self, service, tmp_path
    ):
        spacing = PptxParagraph(
            runs=[PptxTextRun(text="spaced", font_size_pt=16)],
            space_before_pt=3,
            space_after_pt=5,
        )
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    textboxes=[
                        PptxTextBox(
                            left_px=40,
                            top_px=40,
                            width_px=300,
                            height_px=80,
                            paragraphs=[spacing],
                        )
                    ],
                    shapes=[
                        PptxShape(
                            left_px=400,
                            top_px=40,
                            width_px=300,
                            height_px=80,
                            fill_color="#eeeeee",
                            paragraphs=[
                                PptxParagraph(
                                    runs=[
                                        PptxTextRun(
                                            text="shape spaced",
                                            font_size_pt=16,
                                        )
                                    ],
                                    space_before_pt=7,
                                    space_after_pt=9,
                                )
                            ],
                        )
                    ],
                )
            ]
        )

        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        slide = Presentation(response.pptx_path).slides[0]
        paragraphs = {
            shape.text: shape.text_frame.paragraphs[0]
            for shape in slide.shapes
            if shape.has_text_frame
        }
        assert paragraphs["spaced"].space_before.pt == pytest.approx(3)
        assert paragraphs["spaced"].space_after.pt == pytest.approx(5)
        assert paragraphs["shape spaced"].space_before.pt == pytest.approx(7)
        assert paragraphs["shape spaced"].space_after.pt == pytest.approx(9)

    def test_preserves_speaker_notes(self, service, tmp_path):
        spec = _make_design_spec()
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        assert slide.notes_slide.notes_text_frame.text == "발표자 노트"

    def test_preserves_background_color(self, service, tmp_path):
        spec = _make_design_spec()
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        fill = slide.background.fill
        assert fill.fore_color.rgb is not None
        assert str(fill.fore_color.rgb) == "1A1A2E"

    def test_disables_default_theme_effects_for_shapes_and_connectors(
        self, service, tmp_path
    ):
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    shapes=[
                        PptxShape(
                            left_px=40,
                            top_px=40,
                            width_px=240,
                            height_px=100,
                            fill_color="#ffffff",
                            border_color="#000000",
                        ),
                        PptxShape(
                            left_px=300,
                            top_px=90,
                            width_px=180,
                            height_px=0,
                            shape_type="line",
                            border_color="#000000",
                            end_arrow=True,
                        ),
                    ]
                )
            ]
        )

        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        slide = Presentation(response.pptx_path).slides[0]
        effect_refs = []
        for shape in slide.shapes:
            style = shape._element.find(qn("p:style"))
            effect_ref = style.find(qn("a:effectRef")) if style is not None else None
            if effect_ref is not None:
                effect_refs.append(effect_ref.get("idx"))
        assert effect_refs == ["0", "0"]

    def test_explicit_z_index_is_not_overridden_by_textbox_reordering(
        self, service, tmp_path
    ):
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    textboxes=[
                        PptxTextBox(
                            left_px=10,
                            top_px=10,
                            width_px=100,
                            height_px=40,
                            z_index=0,
                            paragraphs=[
                                PptxParagraph(runs=[PptxTextRun(text="behind")])
                            ],
                        )
                    ],
                    shapes=[
                        PptxShape(
                            left_px=10,
                            top_px=10,
                            width_px=100,
                            height_px=40,
                            z_index=1,
                            fill_color="#FF0000",
                        )
                    ],
                )
            ]
        )

        with patch.object(service._builder, "ensure_textboxes_on_top") as ensure:
            response = service.export_from_design_spec(spec, output_dir=tmp_path)

        ensure.assert_not_called()
        slide = Presentation(response.pptx_path).slides[0]
        assert slide.shapes[0].has_text_frame
        assert slide.shapes[0].text == "behind"
        assert slide.shapes[1].shape_type is not None

    def test_multiple_slides(self, service, tmp_path):
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#111111",
                    textboxes=[
                        PptxTextBox(
                            left_px=40,
                            top_px=40,
                            width_px=600,
                            height_px=60,
                            paragraphs=[
                                PptxParagraph(runs=[PptxTextRun(text=f"슬라이드 {i}")])
                            ],
                        )
                    ],
                )
                for i in range(3)
            ]
        )
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        assert len(prs.slides) == 3

    def test_creates_output_dir(self, service, tmp_path):
        nested_dir = tmp_path / "a" / "b" / "export"
        spec = _make_design_spec()
        response = service.export_from_design_spec(spec, output_dir=nested_dir)

        assert nested_dir.exists()
        assert Path(response.pptx_path).exists()

    def test_auto_creates_output_dir(self, service):
        spec = _make_design_spec()
        response = service.export_from_design_spec(spec)

        path = Path(response.pptx_path)
        assert path.exists()
        assert path.name == "presentation.pptx"


class TestBackgroundImage:
    """배경 이미지가 슬라이드 배경 속성(p:bg/a:blipFill)으로 설정되는지 검증."""

    @staticmethod
    def _make_title_spec() -> DesignSpec:
        return DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#1a1a2e",
                    slide_type="title",
                    textboxes=[
                        PptxTextBox(
                            left_px=100,
                            top_px=200,
                            width_px=1080,
                            height_px=80,
                            paragraphs=[
                                PptxParagraph(
                                    runs=[
                                        PptxTextRun(
                                            text="Title",
                                            font_size_pt=40,
                                            bold=True,
                                            color="#ffffff",
                                        )
                                    ]
                                )
                            ],
                        ),
                    ],
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=400,
                            width_px=200,
                            height_px=4,
                            fill_color="#4a90d9",
                            shape_type="rectangle",
                        ),
                    ],
                ),
            ]
        )

    def test_bg_image_set_as_background_property(self, service, tmp_path):
        """배경 이미지는 shape이 아닌 슬라이드 배경 속성(a:blipFill)으로 설정되어야 한다."""
        from io import BytesIO

        from PIL import Image

        buf = BytesIO()
        Image.new("RGB", (4, 4), color="white").save(buf, format="PNG")
        fake_png = buf.getvalue()

        with patch(
            "ppt_generator.tools.pptx.service.bg_image_utils.get_bg_image_bytes",
            return_value=fake_png,
        ):
            spec = self._make_title_spec()
            response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]

        # 배경에 blipFill이 설정되어 있어야 함
        cSld = slide._element.find(qn("p:cSld"))
        bgEl = cSld.find(qn("p:bg"))
        assert bgEl is not None, "p:bg 요소가 존재해야 함"
        bgPr = bgEl.find(qn("p:bgPr"))
        assert bgPr is not None, "p:bgPr 요소가 존재해야 함"
        blipFill = bgPr.find(qn("a:blipFill"))
        assert blipFill is not None, "a:blipFill이 배경에 설정되어야 함"

        # blip의 r:embed가 존재해야 함
        blip = blipFill.find(qn("a:blip"))
        assert blip is not None, "a:blip 요소가 존재해야 함"
        assert blip.get(qn("r:embed")), "r:embed 속성이 있어야 함"

        # shapes에 전체 크기 배경 이미지가 없어야 함 (shape으로 추가되면 안 됨)
        pic_tag = qn("p:pic")
        sp_tree = slide.shapes._spTree
        bg_pics = [
            c
            for c in sp_tree
            if c.tag == pic_tag
            and any(
                s.left == 0 and s.top == 0 and s.width > 12_000_000
                for s in slide.shapes
                if id(s._element) == id(c)
            )
        ]
        assert len(bg_pics) == 0, "배경 이미지가 shape으로 존재하면 안 됨"


class TestConnectorArrow:
    """line shape가 connector로 렌더링되고 화살표 머리가 올바르게 설정되는지 검증."""

    @staticmethod
    def _make_arrow_spec(
        end_arrow: bool = False,
        start_arrow: bool = False,
        dash_style: str | None = None,
    ) -> DesignSpec:
        return DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#232F3E",
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=300,
                            width_px=200,
                            height_px=0,
                            shape_type="line",
                            border_color="#FFC000",
                            border_width_pt=2,
                            end_arrow=end_arrow,
                            start_arrow=start_arrow,
                            dash_style=dash_style,
                        ),
                    ],
                ),
            ]
        )

    def test_line_rendered_as_connector(self, service, tmp_path):
        """line shape는 p:cxnSp(connector)로 렌더링되어야 한다."""
        spec = self._make_arrow_spec()
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        sp_tree = slide.shapes._spTree
        cxn_tag = qn("p:cxnSp")
        connectors = [c for c in sp_tree if c.tag == cxn_tag]
        assert len(connectors) == 1, "line shape가 p:cxnSp(connector)로 렌더링되어야 함"

    def test_end_arrow_triangle(self, service, tmp_path):
        """end_arrow=True → a:tailEnd type=triangle이 있어야 한다."""
        spec = self._make_arrow_spec(end_arrow=True)
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        cxn_tag = qn("p:cxnSp")
        connector = next(c for c in slide.shapes._spTree if c.tag == cxn_tag)

        ln = connector.find(f".//{qn('a:ln')}")
        assert ln is not None
        tail = ln.find(qn("a:tailEnd"))
        assert tail is not None, "end_arrow=True이면 a:tailEnd가 있어야 함"
        assert tail.get("type") == "triangle"

    def test_start_arrow_triangle(self, service, tmp_path):
        """start_arrow=True → a:headEnd type=triangle이 있어야 한다."""
        spec = self._make_arrow_spec(start_arrow=True)
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        cxn_tag = qn("p:cxnSp")
        connector = next(c for c in slide.shapes._spTree if c.tag == cxn_tag)

        ln = connector.find(f".//{qn('a:ln')}")
        assert ln is not None
        head = ln.find(qn("a:headEnd"))
        assert head is not None, "start_arrow=True이면 a:headEnd가 있어야 함"
        assert head.get("type") == "triangle"

    def test_bidirectional_arrows(self, service, tmp_path):
        """양방향 화살표: headEnd와 tailEnd 모두 있어야 한다."""
        spec = self._make_arrow_spec(end_arrow=True, start_arrow=True)
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        cxn_tag = qn("p:cxnSp")
        connector = next(c for c in slide.shapes._spTree if c.tag == cxn_tag)

        ln = connector.find(f".//{qn('a:ln')}")
        assert ln.find(qn("a:tailEnd")) is not None
        assert ln.find(qn("a:headEnd")) is not None

    def test_no_arrow_no_ends(self, service, tmp_path):
        """화살표 미지정 → headEnd/tailEnd가 없어야 한다."""
        spec = self._make_arrow_spec(end_arrow=False, start_arrow=False)
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        cxn_tag = qn("p:cxnSp")
        connector = next(c for c in slide.shapes._spTree if c.tag == cxn_tag)

        ln = connector.find(f".//{qn('a:ln')}")
        assert ln.find(qn("a:tailEnd")) is None
        assert ln.find(qn("a:headEnd")) is None

    def test_dash_style(self, service, tmp_path):
        """dash_style="dash" → a:prstDash val=dash가 있어야 한다."""
        spec = self._make_arrow_spec(dash_style="dash")
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        cxn_tag = qn("p:cxnSp")
        connector = next(c for c in slide.shapes._spTree if c.tag == cxn_tag)

        ln = connector.find(f".//{qn('a:ln')}")
        prstDash = ln.find(qn("a:prstDash"))
        assert prstDash is not None, "dash_style=dash이면 a:prstDash가 있어야 함"
        assert prstDash.get("val") == "dash"

    def test_vertical_connector(self, service, tmp_path):
        """수직 커넥터(width=0, height>0)도 정상 렌더링되어야 한다."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#232F3E",
                    shapes=[
                        PptxShape(
                            left_px=640,
                            top_px=200,
                            width_px=0,
                            height_px=100,
                            shape_type="line",
                            border_color="#FF9900",
                            border_width_pt=2,
                            end_arrow=True,
                        ),
                    ],
                ),
            ]
        )
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        cxn_tag = qn("p:cxnSp")
        connectors = [c for c in slide.shapes._spTree if c.tag == cxn_tag]
        assert len(connectors) == 1

    def test_negative_height_connector_flipV(self, service, tmp_path):
        """음수 height → connector에 flipV="1"이 설정되어야 한다."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#232F3E",
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=100,
                            width_px=200,
                            height_px=-150,
                            shape_type="line",
                            border_color="#FFC000",
                            border_width_pt=2,
                            end_arrow=True,
                        ),
                    ],
                ),
            ]
        )
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        cxn_tag = qn("p:cxnSp")
        connector = next(c for c in slide.shapes._spTree if c.tag == cxn_tag)
        xfrm = connector.find(f".//{qn('a:xfrm')}")
        assert xfrm is not None
        assert xfrm.get("flipV") == "1", "음수 height이면 flipV=1이어야 한다"

    def test_positive_height_no_flipV(self, service, tmp_path):
        """양수 height → connector에 flipV가 없어야 한다."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#232F3E",
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=100,
                            width_px=200,
                            height_px=150,
                            shape_type="line",
                            border_color="#FFC000",
                            border_width_pt=2,
                            end_arrow=True,
                        ),
                    ],
                ),
            ]
        )
        response = service.export_from_design_spec(spec, output_dir=tmp_path)

        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        cxn_tag = qn("p:cxnSp")
        connector = next(c for c in slide.shapes._spTree if c.tag == cxn_tag)
        xfrm = connector.find(f".//{qn('a:xfrm')}")
        assert xfrm is None or xfrm.get("flipV") is None, (
            "양수 height이면 flipV 없어야 한다"
        )


class TestFreeformExport:
    """custGeom freeform export 견고성.

    차트 원호(도넛/파이) 등은 소수 viewBox("403.6 403.6")와 SVG arc("A") 명령을
    쓰는데, 예전엔 int() 파싱과 정수-only tokenizer 라 export 시 crash 했다.
    """

    def test_decimal_viewbox_and_arc_does_not_crash(self, service, tmp_path):
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=100,
                            width_px=200,
                            height_px=200,
                            shape_type="custom",
                            fill_color="#4A90D9",
                            svg_path=(
                                "403.6 403.6 M 201.8 0.0 "
                                "A 201.8 201.8 0 1 1 0.4 189.13 "
                                "L 60.82 192.93 "
                                "A 141.26 141.26 0 1 0 201.8 60.54 Z"
                            ),
                        )
                    ],
                )
            ]
        )
        response = service.export_from_design_spec(
            spec, output_dir=tmp_path, bg_image_policy="none"
        )
        prs = Presentation(response.pptx_path)
        slide = prs.slides[0]
        # freeform sp 가 생성되고 custGeom path 가 정수 좌표를 갖는지 확인
        sp_with_geom = [
            sp
            for sp in slide.shapes._spTree.findall(qn("p:sp"))
            if sp.find(f".//{qn('a:custGeom')}") is not None
        ]
        assert len(sp_with_geom) == 1
        path_el = sp_with_geom[0].find(f".//{qn('a:path')}")
        assert path_el.get("w") == "404" and path_el.get("h") == "404"
        for pt in path_el.findall(f".//{qn('a:pt')}"):
            int(pt.get("x"))  # 정수여야 함 (예외 없이 파싱)
            int(pt.get("y"))

    def test_svg_arcs_export_as_curves_instead_of_chords(self, service, tmp_path):
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=100,
                            width_px=200,
                            height_px=200,
                            shape_type="custom",
                            fill_color="#4A90D9",
                            svg_path=(
                                "400 400 M 200 0 "
                                "A 200 200 0 1 1 0 200 "
                                "L 60 200 "
                                "A 140 140 0 1 0 200 60 Z"
                            ),
                        )
                    ],
                )
            ]
        )

        response = service.export_from_design_spec(
            spec, output_dir=tmp_path, bg_image_policy="none"
        )

        slide = Presentation(response.pptx_path).slides[0]
        freeform = next(
            sp
            for sp in slide.shapes._spTree.findall(qn("p:sp"))
            if sp.find(f".//{qn('a:custGeom')}") is not None
        )
        curves = freeform.findall(f".//{qn('a:cubicBezTo')}")
        assert len(curves) >= 4
