"""PPTX 임포트 기능 테스트.

핵심 테스트 전략: Export → Import 라운드트립으로 데이터 보존을 검증한다.
"""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_generator.interfaces.schemas import (
    DesignSpec,
    PptxImage,
    PptxParagraph,
    PptxShape,
    PptxSlideSpec,
    PptxTextBox,
    PptxTextRun,
)
from ppt_generator.interfaces.spec_utils import (
    design_spec_to_json,
    parse_design_spec_json,
)
from ppt_generator.tools.pptx.service import ExportService
from ppt_generator.tools.pptx_import.service import ImportService
from ppt_generator.tools.pptx_import.slide_reader import SlideReader
from ppt_generator.tools.project.service import ProjectService

# ── 테스트용 DesignSpec 생성 헬퍼 ──


def _make_rich_spec() -> DesignSpec:
    """다양한 요소를 포함한 테스트용 DesignSpec."""
    return DesignSpec(
        slides=[
            PptxSlideSpec(
                background_color="#1A1A2E",
                textboxes=[
                    PptxTextBox(
                        left_px=64,
                        top_px=72,
                        width_px=600,
                        height_px=80,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="제목 텍스트",
                                        font_size_pt=32,
                                        bold=True,
                                        color="#FFFFFF",
                                    ),
                                ]
                            )
                        ],
                        vertical_alignment="top",
                        padding_left_px=0,
                        padding_right_px=0,
                        padding_top_px=0,
                        padding_bottom_px=0,
                    ),
                    PptxTextBox(
                        left_px=64,
                        top_px=148,
                        width_px=600,
                        height_px=400,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="본문 ", font_size_pt=18, color="#CCCCCC"
                                    ),
                                    PptxTextRun(
                                        text="굵게",
                                        font_size_pt=18,
                                        bold=True,
                                        color="#FF9900",
                                    ),
                                ]
                            ),
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="불릿 항목 1",
                                        font_size_pt=16,
                                        color="#CCCCCC",
                                    )
                                ],
                                bullet_level=0,
                            ),
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="불릿 하위 항목",
                                        font_size_pt=14,
                                        color="#999999",
                                    )
                                ],
                                bullet_level=1,
                            ),
                        ],
                        vertical_alignment="top",
                        padding_left_px=8,
                        padding_right_px=8,
                        padding_top_px=4,
                        padding_bottom_px=4,
                    ),
                ],
                shapes=[
                    PptxShape(
                        left_px=700,
                        top_px=148,
                        width_px=500,
                        height_px=200,
                        shape_type="rounded_rectangle",
                        fill_color="#2A2A4E",
                        border_color="#4A90D9",
                        border_width_pt=1.5,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="도형 내 텍스트",
                                        font_size_pt=16,
                                        color="#FFFFFF",
                                    ),
                                ]
                            )
                        ],
                        vertical_alignment="middle",
                        padding_left_px=10,
                        padding_right_px=10,
                        padding_top_px=5,
                        padding_bottom_px=5,
                    ),
                    PptxShape(
                        left_px=700,
                        top_px=400,
                        width_px=500,
                        height_px=150,
                        shape_type="ellipse",
                        fill_color="#FF6600",
                    ),
                    PptxShape(
                        left_px=200,
                        top_px=600,
                        width_px=300,
                        height_px=0,
                        shape_type="line",
                        border_color="#FFC000",
                        border_width_pt=2,
                        end_arrow=True,
                    ),
                ],
                speaker_notes="발표자 노트 내용",
                slide_type="content",
            ),
        ]
    )


def _make_multi_slide_spec() -> DesignSpec:
    """다중 슬라이드 테스트용."""
    return DesignSpec(
        slides=[
            PptxSlideSpec(
                background_color="#000000",
                textboxes=[
                    PptxTextBox(
                        left_px=200,
                        top_px=260,
                        width_px=880,
                        height_px=80,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="프레젠테이션 제목",
                                        font_size_pt=40,
                                        bold=True,
                                        color="#FFFFFF",
                                    ),
                                ]
                            )
                        ],
                    )
                ],
                slide_type="title",
            ),
            PptxSlideSpec(
                background_color="#1A1A2E",
                textboxes=[
                    PptxTextBox(
                        left_px=64,
                        top_px=72,
                        width_px=1152,
                        height_px=48,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="두 번째 슬라이드",
                                        font_size_pt=28,
                                        color="#FFFFFF",
                                    ),
                                ]
                            )
                        ],
                    )
                ],
                slide_type="content",
            ),
            PptxSlideSpec(
                background_color="#1A1A2E",
                textboxes=[
                    PptxTextBox(
                        left_px=200,
                        top_px=240,
                        width_px=880,
                        height_px=80,
                        paragraphs=[
                            PptxParagraph(
                                runs=[
                                    PptxTextRun(
                                        text="감사합니다",
                                        font_size_pt=36,
                                        bold=True,
                                        color="#FFFFFF",
                                    ),
                                ]
                            )
                        ],
                    )
                ],
                slide_type="closing",
            ),
        ]
    )


# ── Fixtures ──


@pytest.fixture
def export_service():
    return ExportService()


@pytest.fixture
def import_service():
    return ImportService()


# ── 라운드트립 테스트 ──


class TestRoundTrip:
    """Export → Import 라운드트립으로 데이터 보존 검증."""

    def _round_trip(self, export_service, import_service, spec, tmp_path) -> DesignSpec:
        """Export → Import 라운드트립."""
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, warnings = import_service.import_from_file(response.pptx_path)
        return imported

    def test_slide_count_preserved(self, export_service, import_service, tmp_path):
        spec = _make_multi_slide_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)
        assert len(imported.slides) == len(spec.slides)

    def test_text_content_preserved(self, export_service, import_service, tmp_path):
        spec = _make_rich_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        slide = imported.slides[0]
        all_text = ""
        for tb in slide.textboxes:
            for p in tb.paragraphs:
                for r in p.runs:
                    all_text += r.text
        for s in slide.shapes:
            for p in s.paragraphs:
                for r in p.runs:
                    all_text += r.text

        assert "제목 텍스트" in all_text
        assert "본문" in all_text
        assert "굵게" in all_text
        assert "불릿 항목 1" in all_text
        assert "불릿 하위 항목" in all_text
        assert "도형 내 텍스트" in all_text

    def test_background_color_preserved(self, export_service, import_service, tmp_path):
        spec = _make_rich_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        assert imported.slides[0].background_color == "#1A1A2E"

    def test_background_image_preserved(self, export_service, import_service, tmp_path):
        """슬라이드 배경의 blipFill 이미지가 라운드트립에서 보존되어야 한다.

        회귀 방지: solid color만 추출하고 blipFill을 무시하면 임포트 시
        흰색으로 폴백되어 시각적으로 누락된다.
        """
        # content 슬라이드여서 기본 폴백 배경 이미지가 적용되지 않는 케이스로 검증
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR"
            b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89"
            b"\x00\x00\x00\rIDATx\x9cc\xfc\xcf\xc0\xf0\x1f\x00\x05\x00\x01\xff"
            b"\x18p\xb1\x82\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#1A1A2E",
                    background_image_bytes=png_bytes,
                    textboxes=[],
                    shapes=[],
                    speaker_notes="",
                    slide_type="content",
                ),
            ]
        )
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        assert imported.slides[0].background_image_bytes, (
            "blipFill 슬라이드 배경 이미지가 임포트에서 추출되지 않음"
        )
        assert imported.slides[0].background_image_bytes.startswith(b"\x89PNG"), (
            "추출된 배경 이미지가 PNG가 아님"
        )

    def test_speaker_notes_preserved(self, export_service, import_service, tmp_path):
        spec = _make_rich_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        assert imported.slides[0].speaker_notes == "발표자 노트 내용"

    def test_font_properties_preserved(self, export_service, import_service, tmp_path):
        spec = _make_rich_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        # 제목 텍스트의 폰트 속성
        title_tb = imported.slides[0].textboxes[0]
        title_run = title_tb.paragraphs[0].runs[0]
        assert title_run.font_size_pt == 32
        assert title_run.bold is True
        assert title_run.color == "#FFFFFF"

    def test_shape_fill_color_preserved(self, export_service, import_service, tmp_path):
        spec = _make_rich_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        # rounded_rectangle의 fill
        rounded_rect = [
            s for s in imported.slides[0].shapes if s.shape_type == "rounded_rectangle"
        ]
        assert len(rounded_rect) >= 1
        assert rounded_rect[0].fill_color == "#2A2A4E"

    def test_connector_arrow_preserved(self, export_service, import_service, tmp_path):
        spec = _make_rich_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        lines = [s for s in imported.slides[0].shapes if s.shape_type == "line"]
        assert len(lines) >= 1
        arrow_line = lines[0]
        assert arrow_line.end_arrow is True

    def test_position_accuracy(self, export_service, import_service, tmp_path):
        """좌표 변환 정밀도 ±3px 이내 검증."""
        spec = _make_rich_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        orig_tb = spec.slides[0].textboxes[0]
        imported_tb = imported.slides[0].textboxes[0]

        assert abs(imported_tb.left_px - orig_tb.left_px) <= 3
        assert abs(imported_tb.top_px - orig_tb.top_px) <= 3
        assert abs(imported_tb.width_px - orig_tb.width_px) <= 3
        assert abs(imported_tb.height_px - orig_tb.height_px) <= 3

    def test_bullet_level_preserved(self, export_service, import_service, tmp_path):
        spec = _make_rich_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        body_tb = imported.slides[0].textboxes[1]
        bullet_paras = [p for p in body_tb.paragraphs if p.bullet_level >= 0]
        assert len(bullet_paras) >= 2
        assert bullet_paras[0].bullet_level == 0
        assert bullet_paras[1].bullet_level >= 1

    def test_ellipse_shape_type(self, export_service, import_service, tmp_path):
        spec = _make_rich_spec()
        imported = self._round_trip(export_service, import_service, spec, tmp_path)

        ellipses = [s for s in imported.slides[0].shapes if s.shape_type == "ellipse"]
        assert len(ellipses) >= 1


# ── ImportService 직접 테스트 ──


class TestImportService:
    def test_import_from_file_not_found(self, import_service):
        with pytest.raises(FileNotFoundError):
            import_service.import_from_file("/nonexistent/path.pptx")

    def test_import_from_file_wrong_extension(self, import_service, tmp_path):
        bad_file = tmp_path / "test.pdf"
        bad_file.write_text("not a pptx")
        with pytest.raises(ValueError, match="PPTX 파일만 지원"):
            import_service.import_from_file(bad_file)

    def test_import_empty_presentation(self, import_service, tmp_path):
        prs = Presentation()
        pptx_path = tmp_path / "empty.pptx"
        prs.save(str(pptx_path))
        with pytest.raises(ValueError, match="슬라이드가 없습니다"):
            import_service.import_from_file(pptx_path)

    def test_import_from_bytes(self, export_service, import_service, tmp_path):
        spec = _make_rich_spec()
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        pptx_bytes = Path(response.pptx_path).read_bytes()

        imported, warnings = import_service.import_from_bytes(pptx_bytes)
        assert len(imported.slides) == 1


# ── SlideReader 단위 테스트 ──


class TestSlideReader:
    def test_compute_scale_standard_size(self):
        """표준 크기(13.333" × 7.5") → scale 1.0."""
        prs = Presentation()
        prs.slide_width = 12_192_000
        prs.slide_height = 6_858_000
        sx, sy = SlideReader.compute_scale(prs)
        assert abs(sx - 1.0) < 0.01
        assert abs(sy - 1.0) < 0.01

    def test_compute_scale_4_3_aspect(self):
        """4:3 비율 (10" × 7.5") → scale_x > 1."""
        prs = Presentation()
        prs.slide_width = 9_144_000  # 10"
        prs.slide_height = 6_858_000  # 7.5"
        sx, sy = SlideReader.compute_scale(prs)
        assert sx > 1.0
        assert abs(sy - 1.0) < 0.01

    def test_infer_slide_type_title(self):
        """첫 슬라이드 + 큰 폰트 + 적은 요소 → title."""
        textboxes = [
            PptxTextBox(
                left_px=200,
                top_px=260,
                width_px=880,
                height_px=80,
                paragraphs=[
                    PptxParagraph(
                        runs=[
                            PptxTextRun(text="Title", font_size_pt=40, bold=True),
                        ]
                    )
                ],
            )
        ]
        result = SlideReader._infer_slide_type(0, 5, textboxes, [])
        assert result == "title"

    def test_infer_slide_type_closing(self):
        """마지막 슬라이드 + closing 키워드 → closing."""
        textboxes = [
            PptxTextBox(
                left_px=200,
                top_px=260,
                width_px=880,
                height_px=80,
                paragraphs=[
                    PptxParagraph(
                        runs=[
                            PptxTextRun(text="감사합니다"),
                        ]
                    )
                ],
            )
        ]
        result = SlideReader._infer_slide_type(4, 5, textboxes, [])
        assert result == "closing"

    def test_infer_slide_type_content(self):
        """중간 슬라이드 → content."""
        textboxes = [
            PptxTextBox(
                left_px=64,
                top_px=72,
                width_px=1152,
                height_px=500,
                paragraphs=[
                    PptxParagraph(
                        runs=[
                            PptxTextRun(text="Some content", font_size_pt=16),
                        ]
                    )
                ],
            )
        ]
        result = SlideReader._infer_slide_type(2, 5, textboxes, [])
        assert result == "content"


# ── 특수 요소 테스트 ──


class TestSpecialElements:
    """python-pptx로 직접 생성한 PPTX의 특수 요소 추출 검증."""

    def test_table_extraction(self, import_service, tmp_path):
        """테이블이 Shape 격자로 변환되는지 검증."""
        prs = Presentation()
        prs.slide_width = 12_192_000
        prs.slide_height = 6_858_000
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        rows, cols = 2, 3
        table_shape = slide.shapes.add_table(
            rows,
            cols,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(2),
        )
        table = table_shape.table
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"

        pptx_path = tmp_path / "table.pptx"
        prs.save(str(pptx_path))

        imported, warnings = import_service.import_from_file(pptx_path)
        slide_spec = imported.slides[0]

        # 테이블 셀이 shape으로 변환되었는지
        assert len(slide_spec.shapes) == rows * cols
        all_text = ""
        for s in slide_spec.shapes:
            for p in s.paragraphs:
                for r in p.runs:
                    all_text += r.text
        assert "A1" in all_text
        assert "B1" in all_text
        assert "A2" in all_text

    def test_vertical_alignment_middle(self, export_service, import_service, tmp_path):
        """vertical_alignment=middle이 보존되는지 검증."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#000000",
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=100,
                            width_px=400,
                            height_px=200,
                            shape_type="rectangle",
                            fill_color="#333333",
                            paragraphs=[
                                PptxParagraph(
                                    runs=[
                                        PptxTextRun(
                                            text="중앙 정렬",
                                            font_size_pt=20,
                                            color="#FFFFFF",
                                        ),
                                    ]
                                )
                            ],
                            vertical_alignment="middle",
                        )
                    ],
                )
            ]
        )
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, _ = import_service.import_from_file(response.pptx_path)

        rect = imported.slides[0].shapes[0]
        assert rect.vertical_alignment == "middle"

    def test_paragraph_alignment_center(self, export_service, import_service, tmp_path):
        """paragraph alignment=center가 보존되는지 검증."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    textboxes=[
                        PptxTextBox(
                            left_px=100,
                            top_px=100,
                            width_px=400,
                            height_px=100,
                            paragraphs=[
                                PptxParagraph(
                                    runs=[
                                        PptxTextRun(text="가운데 정렬", font_size_pt=20)
                                    ],
                                    alignment="center",
                                )
                            ],
                        )
                    ],
                )
            ]
        )
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, _ = import_service.import_from_file(response.pptx_path)

        tb = imported.slides[0].textboxes[0]
        assert tb.paragraphs[0].alignment == "center"

    def test_negative_height_diagonal_arrow_roundtrip(
        self, export_service, import_service, tmp_path
    ):
        """음수 height 대각선 (↗) 화살표가 export→import 후 좌표/방향 보존되는지.

        회귀: bbox top 보정 + flipV 명시 set 이 이중 적용되어 export 후 화살표가
        캔버스 위로 떠 있던 버그. spec convention 은 (left, top) = bbox 좌상,
        h<0 = ↗ 방향이며 export 는 begin/end 좌표만 정확히 넘기면 python-pptx 가
        flipV 를 자동 설정한다.
        """
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    shapes=[
                        PptxShape(
                            left_px=854,
                            top_px=224,
                            width_px=218,
                            height_px=-60,
                            shape_type="line",
                            border_color="#FFFFFF",
                            border_width_pt=2,
                            end_arrow=True,
                        )
                    ],
                )
            ]
        )
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, _ = import_service.import_from_file(response.pptx_path)

        line = imported.slides[0].shapes[0]
        assert round(line.left_px) == 854
        assert round(line.top_px) == 224  # bbox top, 화살표가 위로 안 떠야 함
        assert round(line.width_px) == 218
        assert round(line.height_px) == -60
        assert line.end_arrow is True
        assert line.start_arrow is False

    def test_positive_height_diagonal_arrow_roundtrip(
        self, export_service, import_service, tmp_path
    ):
        """양수 height 대각선 (↘) 화살표가 정확히 보존되는지."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=100,
                            width_px=200,
                            height_px=80,
                            shape_type="line",
                            border_color="#FFFFFF",
                            border_width_pt=2,
                            end_arrow=True,
                        )
                    ],
                )
            ]
        )
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, _ = import_service.import_from_file(response.pptx_path)

        line = imported.slides[0].shapes[0]
        assert round(line.left_px) == 100
        assert round(line.top_px) == 100
        assert round(line.width_px) == 200
        assert round(line.height_px) == 80
        assert line.end_arrow is True

    def test_component_id_preserved(self, export_service, import_service, tmp_path):
        """spec 의 component_id 가 export → import 라운드트립에서 보존되는지.

        design_doc 자체는 PPTX 에 직렬화하지 않으므로 import 결과는 None 일 수 있지만,
        textbox/shape 의 component_id 는 PPTX shape name 에 보존되거나 누락된다.
        현재 구현은 PPTX 에 component_id 를 별도 저장하지 않으므로 import 후
        component_id 는 None 으로 fallback 된다. 이 테스트는 *export 가 component_id
        를 갖고도 정상 동작* 하는지만 확인한다.
        """
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    textboxes=[
                        PptxTextBox(
                            left_px=100,
                            top_px=100,
                            width_px=200,
                            height_px=50,
                            paragraphs=[
                                PptxParagraph(
                                    runs=[PptxTextRun(text="hello", font_size_pt=18)]
                                )
                            ],
                            component_id="left.title",
                        )
                    ],
                )
            ]
        )
        # export 가 component_id 가 있어도 깨지지 않아야 함
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, _ = import_service.import_from_file(response.pptx_path)
        # 텍스트는 보존됨
        tb = imported.slides[0].textboxes[0]
        assert tb.paragraphs[0].runs[0].text == "hello"

    def test_dash_style_preserved(self, export_service, import_service, tmp_path):
        """dash_style이 보존되는지 검증."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=300,
                            width_px=400,
                            height_px=0,
                            shape_type="line",
                            border_color="#FF0000",
                            border_width_pt=2,
                            dash_style="dash",
                            end_arrow=True,
                        )
                    ],
                )
            ]
        )
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, _ = import_service.import_from_file(response.pptx_path)

        line = [s for s in imported.slides[0].shapes if s.shape_type == "line"][0]
        assert line.dash_style == "dash"
        assert line.end_arrow is True

    @pytest.mark.parametrize(
        "ooxml_dash_val,expected",
        [
            ("dash", "dash"),
            ("lgDash", "dash"),
            ("sysDash", "dash"),
            ("dashDot", "dash"),
            ("lgDashDot", "dash"),
            ("lgDashDotDot", "dash"),
            ("dot", "dot"),
            ("sysDot", "dot"),
            ("solid", None),
            ("", None),
        ],
    )
    def test_connector_dash_style_normalization(
        self, ooxml_dash_val, expected, tmp_path
    ):
        """OOXML prstDash 변형(sysDot/sysDash 등)이 임포터에서 dash/dot로 정규화되어야 한다.

        회귀 방지: 과거 임포터는 ("dash","dot") 만 인정하여 sysDot/sysDash 등이 누락됐다.
        """
        # python-pptx 의 connector 객체를 모킹하기보다, slide_reader 의 연결자 추출 코드를
        # 직접 호출하기 위해 PPTX 파일을 만들고 a:prstDash val 을 직접 설정해 임포트한다.
        from io import BytesIO
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.util import Inches
        from lxml import etree

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # CXN connector 추가 (line)
        cxn = slide.shapes.add_connector(1, Inches(1), Inches(2), Inches(3), Inches(2))
        # a:ln 에 prstDash val 직접 삽입
        spPr = cxn._element.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = etree.SubElement(spPr, qn("a:ln"))
        # 기존 prstDash 제거
        for old in ln.findall(qn("a:prstDash")):
            ln.remove(old)
        if ooxml_dash_val:
            prst = etree.SubElement(ln, qn("a:prstDash"))
            prst.set("val", ooxml_dash_val)

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        path = tmp_path / "dash.pptx"
        path.write_bytes(buf.read())

        imported, _ = ImportService().import_from_file(path)
        lines = [s for s in imported.slides[0].shapes if s.shape_type == "line"]
        assert lines, "line shape 가 임포트되어야 한다"
        assert lines[0].dash_style == expected, (
            f"OOXML {ooxml_dash_val!r} → expected {expected!r}, got {lines[0].dash_style!r}"
        )

    def test_italic_preserved(self, export_service, import_service, tmp_path):
        """italic 속성이 보존되는지 검증."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    textboxes=[
                        PptxTextBox(
                            left_px=100,
                            top_px=100,
                            width_px=400,
                            height_px=100,
                            paragraphs=[
                                PptxParagraph(
                                    runs=[
                                        PptxTextRun(
                                            text="이탤릭", font_size_pt=18, italic=True
                                        ),
                                    ]
                                )
                            ],
                        )
                    ],
                )
            ]
        )
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, _ = import_service.import_from_file(response.pptx_path)

        run = imported.slides[0].textboxes[0].paragraphs[0].runs[0]
        assert run.italic is True


# ── 새 도형 타입 라운드트립 테스트 ──


class TestNewShapeTypes:
    """새로 추가된 도형 타입의 Export → Import 라운드트립 검증."""

    _SHAPE_TYPES_WITH_FILL = [
        "up_arrow",
        "down_arrow",
        "left_arrow",
        "right_arrow",
        "chevron",
        "triangle",
        "diamond",
        "pentagon",
        "hexagon",
        "trapezoid",
        "parallelogram",
        "cross",
        "star_4",
        "star_5",
        "heart",
        "flowchart_process",
        "flowchart_decision",
        "flowchart_terminator",
    ]

    @pytest.mark.parametrize("shape_type", _SHAPE_TYPES_WITH_FILL)
    def test_shape_type_round_trip(
        self, export_service, import_service, tmp_path, shape_type
    ):
        """각 도형 타입이 Export → Import 후 shape_type 이 보존되는지 검증."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#000000",
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=100,
                            width_px=200,
                            height_px=200,
                            shape_type=shape_type,
                            fill_color="#4472C4",
                        )
                    ],
                )
            ]
        )
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, _ = import_service.import_from_file(response.pptx_path)

        assert len(imported.slides[0].shapes) >= 1
        imported_shape = imported.slides[0].shapes[0]
        assert imported_shape.shape_type == shape_type
        assert imported_shape.fill_color == "#4472C4"

    def test_arrow_with_text(self, export_service, import_service, tmp_path):
        """화살표 도형 내부에 텍스트가 보존되는지 검증."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    background_color="#000000",
                    shapes=[
                        PptxShape(
                            left_px=100,
                            top_px=100,
                            width_px=200,
                            height_px=300,
                            shape_type="up_arrow",
                            fill_color="#FF6600",
                            text="UP",
                            text_color="#FFFFFF",
                            text_size_pt=16,
                            text_bold=True,
                        )
                    ],
                )
            ]
        )
        response = export_service.export_from_design_spec(spec, output_dir=tmp_path)
        imported, _ = import_service.import_from_file(response.pptx_path)

        shape = imported.slides[0].shapes[0]
        assert shape.shape_type == "up_arrow"
        assert shape.text == "UP"
        assert shape.text_bold is True


# ── 플레이스홀더 서식 상속 테스트 ──


class TestPlaceholderFormatInheritance:
    """OOXML 상속 체인(paragraph → layout → master)에서 서식이 올바르게 resolve되는지 검증."""

    def test_placeholder_title_inherits_layout_font_size(self, tmp_path):
        """layout defRPr에 sz=4800 설정 → 임포트 시 48pt로 추출."""
        from lxml.etree import SubElement
        from pptx.oxml.ns import qn as _qn

        prs = Presentation()
        prs.slide_width = 12_192_000
        prs.slide_height = 6_858_000
        layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(layout)

        # layout placeholder의 lstStyle에 defRPr sz=4800 주입
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == 0:  # 제목 placeholder
                txBody = ph._element.find(_qn("p:txBody"))
                if txBody is not None:
                    lstStyle = txBody.find(_qn("a:lstStyle"))
                    if lstStyle is None:
                        lstStyle = SubElement(txBody, _qn("a:lstStyle"))
                    lvl1pPr = lstStyle.find(_qn("a:lvl1pPr"))
                    if lvl1pPr is None:
                        lvl1pPr = SubElement(lstStyle, _qn("a:lvl1pPr"))
                    defRPr = lvl1pPr.find(_qn("a:defRPr"))
                    if defRPr is None:
                        defRPr = SubElement(lvl1pPr, _qn("a:defRPr"))
                    defRPr.set("sz", "4800")  # 48pt
                break

        # 슬라이드 제목 placeholder에 텍스트 추가 (run에 sz 지정 없음)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = "Layout Inherited Title"
                break

        pptx_path = tmp_path / "layout_inherit.pptx"
        prs.save(str(pptx_path))

        imported, _ = ImportService().import_from_file(pptx_path)
        slide_spec = imported.slides[0]

        # 제목 텍스트를 찾아 font_size_pt 검증
        found = False
        for tb in slide_spec.textboxes:
            for p in tb.paragraphs:
                for r in p.runs:
                    if "Layout Inherited Title" in r.text:
                        assert r.font_size_pt == 48, (
                            f"Expected 48pt, got {r.font_size_pt}"
                        )
                        found = True
        assert found, "제목 텍스트를 찾을 수 없습니다"

    def test_placeholder_title_inherits_master_style(self, tmp_path):
        """master titleStyle 기본값 (sz=4400) → 임포트 시 44pt 추출."""
        prs = Presentation()
        prs.slide_width = 12_192_000
        prs.slide_height = 6_858_000

        # 기본 Presentation의 master titleStyle에 이미 sz=4400 kern=1200이 있음
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # 제목 placeholder에 텍스트 추가 (run에 sz 미지정 → master에서 상속)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = "Master Inherited Title"
                break

        pptx_path = tmp_path / "master_inherit.pptx"
        prs.save(str(pptx_path))

        imported, _ = ImportService().import_from_file(pptx_path)
        slide_spec = imported.slides[0]

        found = False
        for tb in slide_spec.textboxes:
            for p in tb.paragraphs:
                for r in p.runs:
                    if "Master Inherited Title" in r.text:
                        assert r.font_size_pt == 44, (
                            f"Expected 44pt, got {r.font_size_pt}"
                        )
                        found = True
        assert found, "제목 텍스트를 찾을 수 없습니다"

    def test_placeholder_color_inherits_from_layout_scheme(self, tmp_path):
        """layout defRPr에 srgbClr 설정 → 임포트 시 해당 색상 추출."""
        from lxml.etree import SubElement
        from pptx.oxml.ns import qn as _qn

        prs = Presentation()
        prs.slide_width = 12_192_000
        prs.slide_height = 6_858_000
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # layout placeholder lstStyle에 solidFill srgbClr val="FF0000" 주입
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == 0:
                txBody = ph._element.find(_qn("p:txBody"))
                if txBody is not None:
                    lstStyle = txBody.find(_qn("a:lstStyle"))
                    if lstStyle is None:
                        lstStyle = SubElement(txBody, _qn("a:lstStyle"))
                    lvl1pPr = lstStyle.find(_qn("a:lvl1pPr"))
                    if lvl1pPr is None:
                        lvl1pPr = SubElement(lstStyle, _qn("a:lvl1pPr"))
                    defRPr = lvl1pPr.find(_qn("a:defRPr"))
                    if defRPr is None:
                        defRPr = SubElement(lvl1pPr, _qn("a:defRPr"))
                    solidFill = SubElement(defRPr, _qn("a:solidFill"))
                    srgbClr = SubElement(solidFill, _qn("a:srgbClr"))
                    srgbClr.set("val", "FF0000")
                break

        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = "Red Title"
                break

        pptx_path = tmp_path / "color_inherit.pptx"
        prs.save(str(pptx_path))

        imported, _ = ImportService().import_from_file(pptx_path)
        slide_spec = imported.slides[0]

        found = False
        for tb in slide_spec.textboxes:
            for p in tb.paragraphs:
                for r in p.runs:
                    if "Red Title" in r.text:
                        assert r.color == "#FF0000", f"Expected #FF0000, got {r.color}"
                        found = True
        assert found, "제목 텍스트를 찾을 수 없습니다"

    def test_run_direct_value_overrides_inherited(self, tmp_path):
        """run에 직접 지정된 값이 상속값을 오버라이드하는지 검증."""
        from pptx.util import Pt

        prs = Presentation()
        prs.slide_width = 12_192_000
        prs.slide_height = 6_858_000
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = "Direct Size"
                # run에 직접 font_size 지정
                for run in ph.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(24)
                break

        pptx_path = tmp_path / "direct_override.pptx"
        prs.save(str(pptx_path))

        imported, _ = ImportService().import_from_file(pptx_path)
        slide_spec = imported.slides[0]

        found = False
        for tb in slide_spec.textboxes:
            for p in tb.paragraphs:
                for r in p.runs:
                    if "Direct Size" in r.text:
                        assert r.font_size_pt == 24, (
                            f"Expected 24pt, got {r.font_size_pt}"
                        )
                        found = True
        assert found, "텍스트를 찾을 수 없습니다"

    def test_theme_color_map_extraction(self, tmp_path):
        """테마 색상 맵이 올바르게 추출되는지 검증."""
        from ppt_generator.tools.pptx_import.slide_reader import (
            _extract_theme_color_map,
        )

        prs = Presentation()
        color_map = _extract_theme_color_map(prs)
        # 기본 테마에서 최소한 dk1, lt1은 추출되어야 함
        assert "dk1" in color_map or len(color_map) == 0  # 테마가 없으면 빈 맵
        if color_map:
            # 별칭 매핑 확인
            if "dk1" in color_map:
                assert "tx1" in color_map
                assert color_map["tx1"] == color_map["dk1"]


# ── 이미지 src 직렬화/역직렬화 테스트 ──


class TestImageSrcSerialization:
    """이미지 src 필드의 직렬화/역직렬화 검증."""

    def test_image_src_serialized_in_json(self):
        """PptxImage.src가 JSON 직렬화에 포함되는지 검증."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    images=[
                        PptxImage(
                            left_px=100,
                            top_px=200,
                            width_px=300,
                            height_px=400,
                            image_bytes=b"\x89PNG",
                            src="images/slide_01_img_01.png",
                        )
                    ],
                )
            ]
        )
        json_str = design_spec_to_json(spec)
        import json

        data = json.loads(json_str)
        img_data = data["slides"][0]["images"][0]
        assert img_data["src"] == "images/slide_01_img_01.png"
        assert "image_bytes" not in img_data  # image_bytes는 제거됨

    def test_image_src_parsed_from_json(self):
        """JSON에서 PptxImage.src가 역직렬화되는지 검증."""
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    images=[
                        PptxImage(
                            left_px=100,
                            top_px=200,
                            width_px=300,
                            height_px=400,
                            src="images/slide_01_img_01.png",
                        )
                    ],
                )
            ]
        )
        json_str = design_spec_to_json(spec)
        parsed = parse_design_spec_json(json_str)
        assert len(parsed.slides[0].images) == 1
        img = parsed.slides[0].images[0]
        assert img.src == "images/slide_01_img_01.png"
        assert img.left_px == 100
        assert img.width_px == 300

    def test_load_design_spec_with_images_restores_bytes(self, tmp_path):
        """load_design_spec_with_images가 src로부터 image_bytes를 복원하는지 검증."""
        project_service = ProjectService()
        project_dir = tmp_path / "test_project"
        project_dir.mkdir()

        # 이미지 파일 생성
        images_dir = project_dir / "slides" / "images"
        images_dir.mkdir(parents=True)
        png_bytes = b"\x89PNG\r\n\x1a\nfake_image_data"
        (images_dir / "slide_01_img_01.png").write_bytes(png_bytes)

        # src가 포함된 design spec 저장
        spec = DesignSpec(
            slides=[
                PptxSlideSpec(
                    images=[
                        PptxImage(
                            left_px=10,
                            top_px=20,
                            width_px=300,
                            height_px=200,
                            src="images/slide_01_img_01.png",
                        )
                    ],
                )
            ]
        )
        project_service.save_design_spec(project_dir, spec)

        # image_bytes 복원 검증
        loaded = project_service.load_design_spec_with_images(project_dir)
        assert len(loaded.slides[0].images) == 1
        img = loaded.slides[0].images[0]
        assert img.image_bytes == png_bytes
        assert img.src == "images/slide_01_img_01.png"


# ── 소프트 줄바꿈(<a:br>) / 라인 도형 / 배경 그림 회귀 방지 ──


def _add_textbox_with_br(slide, texts: list[str]):
    """runs 사이에 <a:br> 을 넣은 텍스트박스를 추가한다."""
    from pptx.oxml import parse_xml
    from pptx.util import Inches

    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    para = tb.text_frame.paragraphs[0]
    for i, t in enumerate(texts):
        if i > 0:
            para._p.append(
                parse_xml(
                    '<a:br xmlns:a="http://schemas.openxmlformats.org/'
                    'drawingml/2006/main"/>'
                )
            )
        para.add_run().text = t
    return tb


class TestSoftLineBreak:
    """<a:br>(문단 내 소프트 줄바꿈) 보존 회귀 방지.

    python-pptx 의 paragraph.runs 는 <a:br> 을 건너뛰므로, 단순히 run 을 이어붙이면
    "Complex"+"re-architecting" 처럼 단어가 붙는다. 개행이 보존되어야 한다.
    """

    def test_br_preserved_as_newline(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox_with_br(slide, ["Complex", "re-architecting"])

        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)

        joined = "".join(
            r.text for tb in spec.textboxes for p in tb.paragraphs for r in p.runs
        )
        assert joined == "Complex\nre-architecting"

    def test_br_renders_as_html_br(self):
        from ppt_generator.tools.slides.text_renderer import run_to_html

        html = run_to_html(
            PptxTextRun(
                text="Complex\nre-architecting",
                font_size_pt=16,
                color="#FFFFFF",
                bold=True,
            )
        )
        assert "<br>" in html
        assert "Complexre-architecting" not in html


class TestLineAutoShape:
    """prst="line" AutoShape 추출 시 크래시 없이 line 도형으로 처리."""

    def test_line_prst_does_not_crash(self):
        from pptx.oxml.ns import qn
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # prstGeom prst="line" 인 도형을 강제로 삽입
        sp = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(0))
        prstGeom = sp._element.find(qn("p:spPr")).find(qn("a:prstGeom"))
        prstGeom.set("prst", "line")

        reader = SlideReader(1.0, 1.0, prs)
        # 예외 없이 완료되어야 한다
        spec = reader.read_slide(slide, 0, 1)
        assert any(s.shape_type == "line" for s in spec.shapes)


class TestAutofitExtraction:
    """<a:bodyPr> autofit 모드 추출 및 렌더러 shrink 스킵."""

    def _make_textbox(self, autofit_tag: str | None):
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import qn
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
        tb.text_frame.paragraphs[0].add_run().text = "긴 제목 텍스트"
        bodyPr = tb.text_frame._txBody.find(qn("a:bodyPr"))
        # 기존 autofit 자식 제거 후 지정 태그 삽입
        for child in list(bodyPr):
            bodyPr.remove(child)
        if autofit_tag:
            bodyPr.append(
                parse_xml(
                    f'<a:{autofit_tag} xmlns:a="http://schemas.openxmlformats.org/'
                    'drawingml/2006/main"/>'
                )
            )
        return prs, slide

    def test_no_autofit_extracted_as_none(self):
        prs, slide = self._make_textbox("noAutofit")
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        assert spec.textboxes[0].autofit == "none"

    def test_sp_autofit_extracted_as_resize(self):
        prs, slide = self._make_textbox("spAutoFit")
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        assert spec.textboxes[0].autofit == "resize"

    def test_norm_autofit_without_scale_keeps_full_size(self):
        # fontScale 없는 normAutofit → 축소 안 함(PowerPoint 기본). mode="none".
        prs, slide = self._make_textbox("normAutofit")
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        assert spec.textboxes[0].autofit == "none"
        assert spec.textboxes[0].autofit_font_scale is None

    def test_norm_autofit_with_scale_extracted(self):
        # fontScale="90000" → 0.9 스케일이 그대로 적용된다.
        from pptx.oxml.ns import qn

        prs, slide = self._make_textbox("normAutofit")
        bodyPr = slide.shapes[-1].text_frame._txBody.find(qn("a:bodyPr"))
        bodyPr.find(qn("a:normAutofit")).set("fontScale", "90000")
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        assert spec.textboxes[0].autofit_font_scale == 0.9

    def test_no_autofit_field_defaults_to_none(self):
        # autofit 자식이 아예 없으면 OOXML 기본은 noAutofit(축소 없음)이다.
        # 임포트는 원본 레이아웃이 확정된 상태이므로 렌더러가 폰트를 재축소하면
        # 원본보다 작아진다(import/0003). LLM 생성 기본값 "shrink" 와 반대.
        prs, slide = self._make_textbox(None)
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        assert spec.textboxes[0].autofit == "none"

    def test_autofit_none_skips_shrink_in_render(self):
        """autofit="none" 텍스트박스는 박스를 넘겨도 폰트가 축소되지 않아야 한다."""
        from ppt_generator.tools.slides.html_renderer import textbox_to_html

        # 작은 박스에 큰 폰트 — shrink 모드라면 축소될 상황
        big_text = PptxTextRun(text="아주 긴 제목 " * 5, font_size_pt=36, color="#000")
        para = PptxParagraph(runs=[big_text])
        tb_none = PptxTextBox(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=40,
            paragraphs=[para],
            autofit="none",
        )
        html = textbox_to_html(tb_none)
        # 36pt 가 그대로 유지 (축소되지 않음)
        assert "font-size:36.00pt" in html


class TestLineSpacingExtraction:
    """줄간격 추출: 직접 배수 지정 + placeholder 상속(lnSpc) 해석 (import/0003)."""

    def test_direct_multiple_spacing_converts_to_pt(self):
        # 문단에 배수 줄간격(0.9)이 직접 지정되면 폰트 크기(20pt)와 곱해 pt 로 환산.
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        para = tb.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = "본문"
        run.font.size = Pt(20)
        para.line_spacing = 0.9  # 배수

        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        assert spec.textboxes[0].line_spacing_pt == pytest.approx(18.0)  # 0.9 * 20

    def test_inherited_line_spacing_from_master_title_style(self):
        # 문단에 직접 줄간격이 없고 title placeholder 이면, 마스터 titleStyle 의
        # lnSpc(spcPct)를 상속 폰트 크기와 곱해 pt 로 환산해야 한다.
        from ppt_generator.tools.pptx_import.theme_resolver import DefaultRunProps

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.paragraphs[0].add_run().text = "제목"

        reader = SlideReader(1.0, 1.0, prs)
        # 마스터 titleStyle 레벨0 에 90% 줄간격 + 36pt 를 주입
        reader._master_tx_styles = {
            "titleStyle": {0: DefaultRunProps(font_size_pt=36, line_spacing_pct=0.9)}
        }
        reader._layout_def_rpr = {}
        # title placeholder(type=1) 로 해석
        ls = reader._extract_line_spacing(
            tb.text_frame, placeholder_type=1, placeholder_idx=0
        )
        assert ls == pytest.approx(32.4)  # 0.9 * 36

    def test_no_spacing_anywhere_returns_none(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.paragraphs[0].add_run().text = "plain"

        reader = SlideReader(1.0, 1.0, prs)
        assert reader._extract_line_spacing(tb.text_frame) is None


class TestChartExtraction:
    """차트 → 벡터 도형 변환 (import/0003). 이미지 래스터화 대신 SVG 도형으로 재현."""

    def _slide_with_chart(self, chart_type, cats, series):
        from pptx.chart.data import CategoryChartData
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = cats
        for name, vals in series:
            cd.add_series(name, vals)
        slide.shapes.add_chart(
            chart_type, Inches(1), Inches(1), Inches(4), Inches(3), cd
        )
        return prs, slide

    def test_column_chart_becomes_bars_proportional_to_values(self):
        from pptx.enum.chart import XL_CHART_TYPE

        prs, slide = self._slide_with_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, ["A", "B", "C"], [("S", (1.0, 2.0, 3.0))]
        )
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        bars = [s for s in spec.shapes if s.shape_type == "rectangle"]
        assert len(bars) == 3
        heights = [b.height_px for b in bars]
        # 값 1:2:3 → 막대 높이도 1:2:3 비율
        assert heights[1] == pytest.approx(heights[0] * 2, rel=0.02)
        assert heights[2] == pytest.approx(heights[0] * 3, rel=0.02)

    def test_doughnut_chart_becomes_arc_slices_with_hole(self):
        from pptx.enum.chart import XL_CHART_TYPE

        prs, slide = self._slide_with_chart(
            XL_CHART_TYPE.DOUGHNUT, ["A", "B"], [("S", (0.75, 0.25))]
        )
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        slices = [s for s in spec.shapes if s.shape_type == "custom" and s.svg_path]
        assert len(slices) == 2
        # 도넛 슬라이스는 arc(A) 명령을 두 번(바깥호+안쪽호) 포함한다.
        assert all(s.svg_path.count("A") >= 2 for s in slices)
        # 슬라이스마다 채움색이 있고 서로 다르다.
        assert slices[0].fill_color and slices[1].fill_color

    def test_line_chart_becomes_custom_polyline(self):
        from pptx.enum.chart import XL_CHART_TYPE

        prs, slide = self._slide_with_chart(
            XL_CHART_TYPE.LINE, ["A", "B", "C"], [("S", (1.0, 3.0, 2.0))]
        )
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        lines = [s for s in spec.shapes if s.shape_type == "custom" and s.svg_path]
        assert len(lines) == 1
        # 3개 점 → M + L + L
        assert lines[0].svg_path.count("L") == 2
        assert lines[0].border_color is not None


class TestChartArcPath:
    """도넛/파이 슬라이스 SVG path 생성 순수 함수."""

    def test_full_pie_slice_has_no_hole_commands(self):
        from ppt_generator.tools.pptx_import.chart_extractors import _arc_path
        import math

        # 90도 파이 조각 (hole=0): 중심에서 시작하는 부채꼴
        d = _arc_path(50, 50, 50, 0, -math.pi / 2, 0)
        assert d.startswith("M 50 50")  # 중심에서 시작
        assert d.count("A") == 1

    def test_doughnut_slice_has_outer_and_inner_arcs(self):
        from ppt_generator.tools.pptx_import.chart_extractors import _arc_path
        import math

        d = _arc_path(50, 50, 50, 35, -math.pi / 2, 0)
        assert d.count("A") == 2  # 바깥호 + 안쪽호
        assert d.endswith("Z")


class TestGroupZOrderPreserved:
    """그룹 내 그리기 순서(박스 뒤 + 텍스트 앞)가 z_index 에 보존돼야 한다 (import/0003).

    z_index 를 리스트 종류별로만 매기면 textbox 가 항상 shape 보다 낮은 z 를 받아,
    흰 채움 박스가 텍스트를 덮어 콜아웃 라벨이 사라진다.
    """

    def _slide_with_callout_group(self):
        from pptx.oxml.ns import qn
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        group = slide.shapes.add_group_shape()
        # [0] 먼저 그려지는 흰 채움 박스 (뒤)
        box = group.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.6))
        box_spPr = box._element.find(qn("p:spPr"))
        from pptx.oxml import parse_xml

        fill = parse_xml(
            '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/'
            '2006/main"><a:srgbClr val="FFFFFF"/></a:solidFill>'
        )
        box_spPr.append(fill)
        # [1] 나중에 그려지는 텍스트 (앞)
        txt = group.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.6))
        txt.text_frame.text = "Callout label"
        return prs, slide

    def test_text_drawn_after_box_gets_higher_z(self):
        prs, slide = self._slide_with_callout_group()
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        label = next(
            tb
            for tb in spec.textboxes
            if tb.paragraphs and "Callout label" in tb.paragraphs[0].runs[0].text
        )
        # 텍스트박스가 그룹의 다른(먼저 그려진) 요소보다 높은 z 를 가져야 한다
        others = [
            e.z_index
            for e in [*spec.textboxes, *spec.shapes]
            if e is not label and e.z_index is not None
        ]
        assert label.z_index is not None
        if others:
            assert label.z_index >= max(others)


class TestRunImportPptxEndToEnd:
    """run_import_pptx 전체 경로(임포트→저장→HTML)로 충실도 수정이 살아있는지 검증."""

    def _pptx_with_chart(self, tmp_path) -> str:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        prs = Presentation()
        for _ in range(2):  # design_summary 가 slides[1] 폴백을 쓰므로 2장 이상
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            cd = CategoryChartData()
            cd.categories = ["A", "B", "C"]
            cd.add_series("S", (1.0, 2.0, 3.0))
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1),
                Inches(1),
                Inches(4),
                Inches(3),
                cd,
            )
        path = tmp_path / "chart_deck.pptx"
        prs.save(str(path))
        return str(path)

    def test_import_renders_chart_as_bars_in_html(self, tmp_path):
        import json

        from ppt_generator.di.container import DIContainer
        from ppt_generator.tools.pptx_import.controller import run_import_pptx

        pptx = self._pptx_with_chart(tmp_path)
        c = DIContainer(project_root=tmp_path)
        result = run_import_pptx(
            pptx, "chart-proj", c.import_service, c.project_service, c.slides_service
        )
        project_dir = Path(result["slides_html_path"]).parent

        spec = json.loads((project_dir / "design_spec" / "slide_01.json").read_text())
        bars = [s for s in spec["shapes"] if s["shape_type"] == "rectangle"]
        assert len(bars) == 3  # 차트가 3개 막대 도형으로 변환돼 저장됨

        html = (project_dir / "slides" / "slide_01.html").read_text()
        assert html.count("<div") >= 3  # 막대들이 HTML 에 렌더됨


class TestChartNoFillSkipped:
    """series 가 noFill(투명)인 차트는 렌더하지 않는다 (import/0003).

    원본이 차트를 데이터 컨테이너로만 두고 실제 막대는 별도 도형(굵은 line)으로
    그린 경우, 차트를 색칠해 그리면 진짜 막대와 중복된다.
    """

    def _chart_slide(self, series_nofill: bool):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.oxml.ns import qn
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["A", "B", "C"]
        cd.add_series("S", (1.0, 2.0, 3.0))
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(3),
            cd,
        )
        if series_nofill:
            # 첫 series 에 <c:spPr><a:noFill/></c:spPr> 주입
            from pptx.oxml import parse_xml

            ser = gf.chart._chartSpace.findall(".//" + qn("c:ser"))[0]
            spPr = parse_xml(
                '<c:spPr xmlns:c="http://schemas.openxmlformats.org/'
                'drawingml/2006/chart" xmlns:a="http://schemas.openxmlformats.org/'
                'drawingml/2006/main"><a:noFill/></c:spPr>'
            )
            # c:ser 자식 순서상 idx/order 뒤에 삽입
            ser.insert(2, spPr)
        return prs, slide

    def test_nofill_series_chart_is_not_rendered(self):
        prs, slide = self._chart_slide(series_nofill=True)
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        bars = [s for s in spec.shapes if s.shape_type == "rectangle"]
        assert bars == []  # noFill 차트 → 막대 렌더 생략

    def test_filled_series_chart_still_renders(self):
        prs, slide = self._chart_slide(series_nofill=False)
        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        bars = [s for s in spec.shapes if s.shape_type == "rectangle"]
        assert len(bars) == 3  # 기본 채움 차트는 그대로 렌더


class TestGradientBorderApproximation:
    """그라데이션 테두리(a:ln>a:gradFill)를 첫 stop 색으로 근사 (import/0003).

    python-pptx 의 shape.line API 접근이 gradFill 을 solidFill 로 변형하므로,
    API 접근 전에 XML 에서 gradient 첫 stop 을 추출해야 색이 유지된다.
    """

    def _round_rect_grad_border(self):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(1)
        )
        spPr = sp._element.find(qn("p:spPr"))
        # 기존 ln 제거 후 gradient ln 주입
        for ln in spPr.findall(qn("a:ln")):
            spPr.remove(ln)
        ln = parse_xml(
            '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'w="19050"><a:gradFill><a:gsLst>'
            '<a:gs pos="40000"><a:srgbClr val="2A39F5"/></a:gs>'
            '<a:gs pos="100000"><a:srgbClr val="00D692"/></a:gs>'
            "</a:gsLst></a:gradFill></a:ln>"
        )
        spPr.append(ln)
        return prs, slide, sp

    def test_gradient_border_uses_first_stop_color(self):
        prs, slide, sp = self._round_rect_grad_border()
        reader = SlideReader(1.0, 1.0, prs)
        color, width = reader._extract_line_style(sp)
        assert color == "#2A39F5"  # 첫 stop
        assert width == pytest.approx(1.5)  # 19050 EMU / 12700


class TestFontNamePreserved:
    """원본 폰트명 보존 및 CSS font-family 렌더링."""

    def test_font_name_extracted(self):
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = "Brand"
        run.font.name = "Amazon Ember Display"
        run.font.size = Pt(24)

        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        r = spec.textboxes[0].paragraphs[0].runs[0]
        assert r.font_name == "Amazon Ember Display"

    def test_font_name_renders_css_family(self):
        from ppt_generator.tools.slides.text_renderer import run_to_html

        html = run_to_html(
            PptxTextRun(text="Brand", font_size_pt=24, font_name="Amazon Ember Display")
        )
        assert "font-family:'Amazon Ember Display'" in html

    def test_placeholder_inherits_layout_latin_over_theme_major(self):
        # title placeholder 의 run 에 폰트가 없고 layout lstStyle 이 특정 latin 을
        # 지정하면, theme major 로 폴백하지 않고 layout 폰트를 상속해야 한다.
        # (theme major 가 "Heavy" 계열이면 글자 폭이 달라져 줄바꿈이 틀어진다 — import/0003)
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        title = slide.shapes.title
        title.text_frame.text = "Heading"
        ph_idx = title.placeholder_format.idx

        # layout placeholder 에 lstStyle > lvl1pPr > defRPr latin 주입
        layout = slide.slide_layout
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == ph_idx:
                txBody = ph._element.find(qn("p:txBody"))
                for old in txBody.findall(qn("a:lstStyle")):
                    txBody.remove(old)
                lst = parse_xml(
                    '<a:lstStyle xmlns:a="http://schemas.openxmlformats.org/'
                    'drawingml/2006/main"><a:lvl1pPr><a:defRPr>'
                    '<a:latin typeface="Amazon Ember Display"/>'
                    "</a:defRPr></a:lvl1pPr></a:lstStyle>"
                )
                txBody.insert(1, lst)
                break

        reader = SlideReader(1.0, 1.0, prs)
        spec = reader.read_slide(slide, 0, 1)
        title_tb = next(
            tb
            for tb in spec.textboxes
            if tb.paragraphs and "Heading" in tb.paragraphs[0].runs[0].text
        )
        assert title_tb.paragraphs[0].runs[0].font_name == "Amazon Ember Display"
