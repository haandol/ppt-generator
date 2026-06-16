"""Tests for shape_builders: padding and vertical_alignment."""

from __future__ import annotations

from pptx.oxml.ns import qn
from pptx.util import Emu

from ppt_generator.interfaces.constants import (
    PPTX_SHAPE_DEFAULT_MARGIN_LR_EMU,
    PPTX_SHAPE_DEFAULT_MARGIN_TB_EMU,
    PX_TO_EMU,
)
from ppt_generator.interfaces.schemas import PptxParagraph, PptxShape, PptxTextRun
from ppt_generator.tools.pptx.shape_builders import add_auto_shape_from_spec


# ── padding: text 경로 ─────────────────────────────────────────


class TestTextPathPadding:
    """shape_spec.text 경로에서 padding 적용 검증."""

    def test_default_padding_when_none(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(left_px=0, top_px=0, width_px=200, height_px=100, text="hello")
        add_auto_shape_from_spec(slide, spec)

        tf = slide.shapes[0].text_frame
        assert tf.margin_left == Emu(PPTX_SHAPE_DEFAULT_MARGIN_LR_EMU)
        assert tf.margin_right == Emu(PPTX_SHAPE_DEFAULT_MARGIN_LR_EMU)
        assert tf.margin_top == Emu(PPTX_SHAPE_DEFAULT_MARGIN_TB_EMU)
        assert tf.margin_bottom == Emu(PPTX_SHAPE_DEFAULT_MARGIN_TB_EMU)

    def test_custom_padding(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=100,
            text="hello",
            padding_left_px=20,
            padding_right_px=10,
            padding_top_px=15,
            padding_bottom_px=5,
        )
        add_auto_shape_from_spec(slide, spec)

        tf = slide.shapes[0].text_frame
        assert tf.margin_left == Emu(int(20 * PX_TO_EMU))
        assert tf.margin_right == Emu(int(10 * PX_TO_EMU))
        assert tf.margin_top == Emu(int(15 * PX_TO_EMU))
        assert tf.margin_bottom == Emu(int(5 * PX_TO_EMU))

    def test_partial_padding_uses_default_for_missing(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=100,
            text="hello",
            padding_left_px=30,
        )
        add_auto_shape_from_spec(slide, spec)

        tf = slide.shapes[0].text_frame
        assert tf.margin_left == Emu(int(30 * PX_TO_EMU))
        assert tf.margin_right == Emu(PPTX_SHAPE_DEFAULT_MARGIN_LR_EMU)
        assert tf.margin_top == Emu(PPTX_SHAPE_DEFAULT_MARGIN_TB_EMU)
        assert tf.margin_bottom == Emu(PPTX_SHAPE_DEFAULT_MARGIN_TB_EMU)


# ── padding: paragraphs 경로 ───────────────────────────────────


class TestParagraphsPathPadding:
    """shape_spec.paragraphs 경로에서 padding 적용 검증."""

    def _make_paragraph(self, text="test"):
        return PptxParagraph(runs=[PptxTextRun(text=text)])

    def test_default_padding_when_none(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=100,
            paragraphs=[self._make_paragraph()],
        )
        add_auto_shape_from_spec(slide, spec)

        tf = slide.shapes[0].text_frame
        assert tf.margin_left == Emu(PPTX_SHAPE_DEFAULT_MARGIN_LR_EMU)
        assert tf.margin_right == Emu(PPTX_SHAPE_DEFAULT_MARGIN_LR_EMU)
        assert tf.margin_top == Emu(PPTX_SHAPE_DEFAULT_MARGIN_TB_EMU)
        assert tf.margin_bottom == Emu(PPTX_SHAPE_DEFAULT_MARGIN_TB_EMU)

    def test_custom_padding(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=100,
            paragraphs=[self._make_paragraph()],
            padding_left_px=16,
            padding_right_px=16,
            padding_top_px=14,
            padding_bottom_px=14,
        )
        add_auto_shape_from_spec(slide, spec)

        tf = slide.shapes[0].text_frame
        assert tf.margin_left == Emu(int(16 * PX_TO_EMU))
        assert tf.margin_right == Emu(int(16 * PX_TO_EMU))
        assert tf.margin_top == Emu(int(14 * PX_TO_EMU))
        assert tf.margin_bottom == Emu(int(14 * PX_TO_EMU))


# ── vertical_alignment: text 경로 ─────────────────────────────


def _get_anchor(text_frame):
    bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
    return bodyPr.get("anchor") if bodyPr is not None else None


class TestTextPathVerticalAlignment:
    """shape_spec.text 경로에서 vertical_alignment 검증."""

    def test_default_is_top(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(left_px=0, top_px=0, width_px=200, height_px=100, text="hello")
        add_auto_shape_from_spec(slide, spec)

        assert _get_anchor(slide.shapes[0].text_frame) == "t"

    def test_middle_alignment(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=100,
            text="hello",
            vertical_alignment="middle",
        )
        add_auto_shape_from_spec(slide, spec)

        assert _get_anchor(slide.shapes[0].text_frame) == "ctr"

    def test_bottom_alignment(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=100,
            text="hello",
            vertical_alignment="bottom",
        )
        add_auto_shape_from_spec(slide, spec)

        assert _get_anchor(slide.shapes[0].text_frame) == "b"


# ── vertical_alignment: paragraphs 경로 ───────────────────────


class TestParagraphsPathVerticalAlignment:
    """shape_spec.paragraphs 경로에서 vertical_alignment 검증."""

    def _make_paragraph(self, text="test"):
        return PptxParagraph(runs=[PptxTextRun(text=text)])

    def test_default_is_top(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=100,
            paragraphs=[self._make_paragraph()],
        )
        add_auto_shape_from_spec(slide, spec)

        assert _get_anchor(slide.shapes[0].text_frame) == "t"

    def test_middle_alignment(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=100,
            paragraphs=[self._make_paragraph()],
            vertical_alignment="middle",
        )
        add_auto_shape_from_spec(slide, spec)

        assert _get_anchor(slide.shapes[0].text_frame) == "ctr"

    def test_bottom_alignment(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=200,
            height_px=100,
            paragraphs=[self._make_paragraph()],
            vertical_alignment="bottom",
        )
        add_auto_shape_from_spec(slide, spec)

        assert _get_anchor(slide.shapes[0].text_frame) == "b"


# ── 컴팩트 라벨(번호 뱃지) word_wrap 억제 ─────────────────────────


class TestCompactLabelNoWrap:
    """작은 도형 + 짧은 텍스트(번호 뱃지)는 PPTX 줄바꿈을 끄고 여백 0."""

    def test_compact_badge_text_disables_wrap(self, blank_slide):
        slide = blank_slide
        # 36x36 원형에 "01" — 전형적 번호 뱃지.
        spec = PptxShape(
            left_px=82,
            top_px=166,
            width_px=36,
            height_px=36,
            shape_type="ellipse",
            text="01",
            text_size_pt=14,
            text_bold=True,
        )
        add_auto_shape_from_spec(slide, spec)

        tf = slide.shapes[0].text_frame
        assert tf.word_wrap is False
        assert tf.margin_left == Emu(0)
        assert tf.margin_right == Emu(0)

    def test_compact_badge_paragraphs_disables_wrap(self, blank_slide):
        slide = blank_slide
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=40,
            height_px=40,
            shape_type="ellipse",
            paragraphs=[
                PptxParagraph(runs=[PptxTextRun(text="02", font_size_pt=14, bold=True)])
            ],
        )
        add_auto_shape_from_spec(slide, spec)

        tf = slide.shapes[0].text_frame
        assert tf.word_wrap is False

    def test_large_card_keeps_wrap(self, blank_slide):
        slide = blank_slide
        # 큰 카드 + 긴 본문 → 줄바꿈 유지(과교정 방지).
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=400,
            height_px=200,
            shape_type="rounded_rectangle",
            text="이것은 여러 줄로 줄바꿈되어야 하는 긴 카드 본문 텍스트입니다.",
        )
        add_auto_shape_from_spec(slide, spec)

        tf = slide.shapes[0].text_frame
        assert tf.word_wrap is True

    def test_small_shape_long_text_keeps_wrap(self, blank_slide):
        slide = blank_slide
        # 작아도 텍스트가 길면 라벨이 아니므로 줄바꿈 유지.
        spec = PptxShape(
            left_px=0,
            top_px=0,
            width_px=40,
            height_px=40,
            shape_type="ellipse",
            text="긴 텍스트입니다",
        )
        add_auto_shape_from_spec(slide, spec)

        tf = slide.shapes[0].text_frame
        assert tf.word_wrap is True
