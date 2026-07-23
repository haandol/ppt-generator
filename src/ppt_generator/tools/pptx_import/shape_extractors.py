"""개별 도형 유형별 추출 모듈.

SlideReader에서 사용하는 각 shape 유형(AutoShape, Freeform, Connector,
Image, TextBox)의 추출 로직을 담당한다.
Group/Table은 compound_extractors 모듈에서 처리한다.
"""

from __future__ import annotations

import logging
import re
from io import BytesIO
from xml.etree.ElementTree import fromstring

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from ppt_generator.interfaces.schemas import (
    PptxImage,
    PptxParagraph,
    PptxShape,
    PptxTextBox,
)
from ppt_generator.tools.pptx_import.chart_extractors import ChartExtractorMixin
from ppt_generator.tools.pptx_import.compound_extractors import CompoundExtractorMixin
from ppt_generator.tools.pptx_import.ooxml_utils import (
    SHAPE_TYPE_REVERSE_MAP,
    custgeom_to_svg_path,
)

logger = logging.getLogger(__name__)

_SVG_BLIP_TAG = "{http://schemas.microsoft.com/office/drawing/2016/SVG/main}svgBlip"
_THEME_FILL_CLASS_RE = re.compile(
    r"^MsftOfcThm_(Text|Background|Accent)([1-6])_Fill(?:_v\d+)?$"
)
_SVG_GRAPHIC_TAGS = frozenset(
    {"path", "rect", "circle", "ellipse", "polygon", "polyline", "line"}
)


def _svg_theme_fill_key(svg_bytes: bytes) -> str | None:
    """SVG 도형이 단일 Office 테마 fill class를 쓰면 scheme key를 반환."""
    try:
        root = fromstring(svg_bytes)
    except Exception:
        return None

    keys: set[str] = set()
    for element in root.iter():
        local_name = element.tag.rsplit("}", 1)[-1]
        if local_name not in _SVG_GRAPHIC_TAGS or element.get("fill") == "none":
            continue
        for class_name in element.get("class", "").split():
            match = _THEME_FILL_CLASS_RE.match(class_name)
            if match is None:
                continue
            family, number = match.groups()
            prefix = {
                "Text": "tx",
                "Background": "bg",
                "Accent": "accent",
            }[family]
            keys.add(f"{prefix}{number}")
    return next(iter(keys)) if len(keys) == 1 else None


def _recolor_svg_theme_fallback(
    fallback_bytes: bytes,
    svg_bytes: bytes,
    theme_color_map: dict[str, str],
) -> bytes:
    """테마 단색 SVG의 PNG fallback을 현재 슬라이드 테마색으로 재색칠."""
    theme_key = _svg_theme_fill_key(svg_bytes)
    color = theme_color_map.get(theme_key) if theme_key else None
    if not color or not re.fullmatch(r"#[0-9A-Fa-f]{6}", color):
        return fallback_bytes

    try:
        from PIL import Image

        with Image.open(BytesIO(fallback_bytes)) as image:
            rgba = image.convert("RGBA")
        red = int(color[1:3], 16)
        green = int(color[3:5], 16)
        blue = int(color[5:7], 16)
        alpha = rgba.getchannel("A")
        recolored = Image.new("RGBA", rgba.size, (red, green, blue, 0))
        recolored.putalpha(alpha)
        output = BytesIO()
        recolored.save(output, format="PNG")
        return output.getvalue()
    except Exception:
        logger.debug("SVG 테마 fallback 재색칠 실패", exc_info=True)
        return fallback_bytes


def _extract_svg_blip_bytes(shape) -> bytes:
    """picture의 asvg:svgBlip 관계에서 원본 SVG bytes를 추출."""
    try:
        svg_blip = shape._element.find(f".//{_SVG_BLIP_TAG}")
        if svg_blip is None:
            return b""
        rel_id = svg_blip.get(qn("r:embed"))
        if not rel_id:
            return b""
        return shape.part.rels[rel_id].target_part.blob
    except Exception:
        logger.debug(
            "SVG picture 관계 추출 실패 (name=%s)",
            getattr(shape, "name", "?"),
            exc_info=True,
        )
        return b""


class ShapeExtractorMixin(CompoundExtractorMixin, ChartExtractorMixin):
    """도형별 추출 기능을 제공하는 mixin 클래스.

    SlideReader에서 상속하여 사용한다.
    TextExtractorMixin, StyleExtractorMixin, 좌표 변환 메서드가 필요하다.
    """

    # 타입 힌트 (실제 구현은 SlideReader/다른 mixin에서 제공)
    # _emu_to_px_x, _emu_to_px_y, _extract_paragraphs, _extract_shape는
    # 부모 CompoundExtractorMixin에서 선언
    _emu_to_px_padding: any
    _extract_line_spacing: any
    _line_spacing_is_default: any
    _extract_vertical_alignment: any
    _extract_fill_color: any
    _extract_fill_gradient: any
    _extract_line_style: any
    _extract_dash_style_from_shape: any

    # ── Shape 분기 ──

    @staticmethod
    def _get_placeholder_info(shape) -> tuple[int | None, int | None]:
        """shape에서 placeholder type과 idx를 추출."""
        try:
            ph_fmt = shape.placeholder_format
        except (ValueError, AttributeError):
            return None, None
        if ph_fmt is None:
            return None, None
        ph_type = getattr(ph_fmt, "type", None)
        ph_idx = getattr(ph_fmt, "idx", None)
        return (int(ph_type) if ph_type is not None else None, ph_idx)

    def _extract_shape(
        self,
        shape,
        textboxes: list[PptxTextBox],
        shapes: list[PptxShape],
        images: list[PptxImage],
        warnings: list[str],
    ) -> None:
        shape_type = getattr(shape, "shape_type", None)

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            self._extract_group(shape, textboxes, shapes, images, warnings)
            return

        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append(self._extract_image(shape))
            return

        if shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                try:
                    _ = shape.image.blob
                    images.append(self._extract_image(shape))
                    return
                except Exception:
                    logger.debug(
                        "플레이스홀더 이미지 추출 실패 (name=%s)",
                        getattr(shape, "name", "?"),
                        exc_info=True,
                    )

        if shape_type == MSO_SHAPE_TYPE.TABLE:
            self._extract_table(shape, shapes, warnings)
            return

        el_tag = shape._element.tag if hasattr(shape, "_element") else ""
        if el_tag == qn("p:cxnSp"):
            shapes.append(self._extract_connector(shape))
            return

        if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            tb = self._extract_textbox(shape)
            if tb.paragraphs:
                textboxes.append(tb)
            return

        if shape_type == MSO_SHAPE_TYPE.FREEFORM:
            spPr = shape._element.find(qn("p:spPr"))
            svg_path = custgeom_to_svg_path(spPr) if spPr is not None else None
            if svg_path is not None:
                shapes.append(self._extract_freeform_shape(shape, svg_path))
                return

        if shape_type in (
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
        ):
            shapes.append(self._extract_auto_shape(shape))
            return

        if shape_type == MSO_SHAPE_TYPE.CHART:
            # 차트를 이미지 대신 벡터 도형으로 재현 (import/0003). 미지원 유형이면
            # _extract_chart 가 경고를 남기고 False 를 반환하므로 그대로 스킵한다.
            self._extract_chart(shape, textboxes, shapes, warnings)
            return

        if shape_type == MSO_SHAPE_TYPE.MEDIA:
            warnings.append(f"미디어 요소는 미지원입니다 (name={shape.name})")
            return

        if hasattr(shape, "text_frame"):
            ph_type, ph_idx = self._get_placeholder_info(shape)
            tf = shape.text_frame
            paragraphs = self._extract_paragraphs(
                tf,
                placeholder_type=ph_type,
                placeholder_idx=ph_idx,
            )
            if paragraphs:
                ph_autofit_mode, ph_autofit_scale = self._extract_autofit(tf)
                textboxes.append(
                    PptxTextBox(
                        left_px=self._emu_to_px_x(shape.left),
                        top_px=self._emu_to_px_y(shape.top),
                        width_px=self._emu_to_px_x(shape.width),
                        height_px=self._emu_to_px_y(shape.height),
                        paragraphs=paragraphs,
                        line_spacing_pt=self._extract_line_spacing(
                            tf,
                            placeholder_type=ph_type,
                            placeholder_idx=ph_idx,
                        ),
                        line_spacing_is_default=self._line_spacing_is_default(
                            tf,
                            placeholder_type=ph_type,
                            placeholder_idx=ph_idx,
                        ),
                        vertical_alignment=self._extract_vertical_alignment(
                            tf, placeholder_idx=ph_idx
                        ),
                        autofit=ph_autofit_mode,
                        autofit_font_scale=ph_autofit_scale,
                        padding_left_px=self._emu_to_px_padding(tf.margin_left),
                        padding_right_px=self._emu_to_px_padding(tf.margin_right),
                        padding_top_px=self._emu_to_px_padding(tf.margin_top),
                        padding_bottom_px=self._emu_to_px_padding(tf.margin_bottom),
                    )
                )

    # ── 텍스트박스 추출 ──

    def _extract_textbox(self, shape) -> PptxTextBox:
        tf = shape.text_frame
        paragraphs = self._extract_paragraphs(tf)
        line_spacing = self._extract_line_spacing(tf)
        vertical_alignment = self._extract_vertical_alignment(tf)
        autofit_mode, autofit_scale = self._extract_autofit(tf)

        return PptxTextBox(
            left_px=self._emu_to_px_x(shape.left),
            top_px=self._emu_to_px_y(shape.top),
            width_px=self._emu_to_px_x(shape.width),
            height_px=self._emu_to_px_y(shape.height),
            paragraphs=paragraphs,
            line_spacing_pt=line_spacing,
            line_spacing_is_default=self._line_spacing_is_default(tf),
            vertical_alignment=vertical_alignment,
            autofit=autofit_mode,
            autofit_font_scale=autofit_scale,
            padding_left_px=self._emu_to_px_padding(tf.margin_left),
            padding_right_px=self._emu_to_px_padding(tf.margin_right),
            padding_top_px=self._emu_to_px_padding(tf.margin_top),
            padding_bottom_px=self._emu_to_px_padding(tf.margin_bottom),
        )

    # ── Shape 텍스트 공통 추출 헬퍼 ──

    def _extract_shape_text_props(self, shape) -> dict:
        """AutoShape/Freeform 공통: text_frame에서 텍스트 관련 속성을 추출."""
        paragraphs: list[PptxParagraph] = []
        line_spacing: float | None = None
        line_spacing_is_default = False
        vertical_alignment = "top"
        padding_left: float | None = None
        padding_right: float | None = None
        padding_top: float | None = None
        padding_bottom: float | None = None
        text: str | None = None
        text_color: str | None = None
        text_size_pt: int | None = None
        text_bold = False

        if shape.has_text_frame:
            tf = shape.text_frame
            paragraphs = self._extract_paragraphs(tf)
            line_spacing = self._extract_line_spacing(tf)
            line_spacing_is_default = self._line_spacing_is_default(tf)
            vertical_alignment = self._extract_vertical_alignment(tf)
            padding_left = self._emu_to_px_padding(tf.margin_left)
            padding_right = self._emu_to_px_padding(tf.margin_right)
            padding_top = self._emu_to_px_padding(tf.margin_top)
            padding_bottom = self._emu_to_px_padding(tf.margin_bottom)

            if (
                len(paragraphs) == 1
                and len(paragraphs[0].runs) == 1
                and paragraphs[0].bullet_level < 0
            ):
                run = paragraphs[0].runs[0]
                text = run.text
                text_color = run.color
                text_size_pt = run.font_size_pt
                text_bold = run.bold

        return {
            "paragraphs": paragraphs,
            "line_spacing_pt": line_spacing,
            "line_spacing_is_default": line_spacing_is_default,
            "vertical_alignment": vertical_alignment,
            "padding_left_px": padding_left,
            "padding_right_px": padding_right,
            "padding_top_px": padding_top,
            "padding_bottom_px": padding_bottom,
            "text": text,
            "text_color": text_color,
            "text_size_pt": text_size_pt,
            "text_bold": text_bold,
        }

    # ── Freeform (custGeom) 추출 ──

    def _extract_freeform_shape(self, shape, svg_path: str) -> PptxShape:
        """Freeform custGeom 도형 → PptxShape (shape_type="custom", svg_path 포함)."""
        fill_color = self._extract_fill_color(shape)
        border_color, border_width = self._extract_line_style(shape)
        text_props = self._extract_shape_text_props(shape)

        return PptxShape(
            left_px=self._emu_to_px_x(shape.left),
            top_px=self._emu_to_px_y(shape.top),
            width_px=self._emu_to_px_x(shape.width),
            height_px=self._emu_to_px_y(shape.height),
            shape_type="custom",
            fill_color=fill_color,
            border_color=border_color,
            border_width_pt=border_width,
            svg_path=svg_path,
            rotation=self._read_rotation(shape),
            **text_props,
        )

    # ── AutoShape 추출 ──

    @staticmethod
    def _read_rotation(shape) -> float:
        """spPr > xfrm@rot(60000분의 1도)을 시계방향 degree 로 읽는다. 없으면 0."""
        try:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is None:
                return 0.0
            xfrm = spPr.find(qn("a:xfrm"))
            if xfrm is None:
                return 0.0
            rot = xfrm.get("rot")
            if rot is None:
                return 0.0
            # OOXML rot 은 60000 단위 = 1도. 0~360 으로 정규화.
            return (int(rot) / 60000.0) % 360.0
        except (ValueError, TypeError):
            return 0.0

    @staticmethod
    def _read_prst(shape) -> str | None:
        """spPr > prstGeom 의 prst 속성값을 직접 읽는다 (enum 매핑 우회)."""
        try:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is None:
                return None
            prstGeom = spPr.find(qn("a:prstGeom"))
            if prstGeom is None:
                return None
            return prstGeom.get("prst")
        except Exception:
            return None

    def _extract_auto_shape(self, shape) -> PptxShape:
        # auto_shape_type 접근은 prst 값이 MSO_AUTO_SHAPE_TYPE 에 없는 경우
        # (예: prst="line") ValueError 를 던진다. prst 를 직접 읽어 매핑을 시도하고,
        # 실패하면 rectangle 로 폴백한다.
        try:
            auto_shape_type = shape.auto_shape_type
        except (ValueError, TypeError, AttributeError):
            auto_shape_type = None

        type_str = SHAPE_TYPE_REVERSE_MAP.get(auto_shape_type)
        if type_str is None:
            prst = self._read_prst(shape)
            if prst == "line":
                # prst="line" 은 커넥터가 아닌 선 도형 → line 으로 처리
                return self._extract_connector(shape)
            type_str = "rectangle"

        fill_gradient = self._extract_fill_gradient(shape)
        fill_color = (
            fill_gradient.stops[0].color
            if fill_gradient is not None
            else self._extract_fill_color(shape)
        )
        border_color, border_width = self._extract_line_style(shape)
        dash_style = self._extract_dash_style_from_shape(shape)
        text_props = self._extract_shape_text_props(shape)

        return PptxShape(
            left_px=self._emu_to_px_x(shape.left),
            top_px=self._emu_to_px_y(shape.top),
            width_px=self._emu_to_px_x(shape.width),
            height_px=self._emu_to_px_y(shape.height),
            shape_type=type_str,
            fill_color=fill_color,
            fill_gradient=fill_gradient,
            border_color=border_color,
            border_width_pt=border_width,
            dash_style=dash_style,
            rotation=self._read_rotation(shape),
            **text_props,
        )

    # ── 커넥터 (선/화살표) 추출 ──

    def _extract_connector(self, shape) -> PptxShape:
        el = shape._element

        left_emu = shape.left
        top_emu = shape.top
        width_emu = shape.width
        height_emu = shape.height

        border_color, border_width = self._extract_line_style(shape)

        spPr = el.find(qn("p:spPr"))

        xfrm = spPr.find(qn("a:xfrm")) if spPr is not None else None
        flip_h = xfrm is not None and xfrm.get("flipH") == "1"
        flip_v = xfrm is not None and xfrm.get("flipV") == "1"

        # bentConnector(꺾인 커넥터) 감지: 직선 대신 직각(elbow) 폴리라인으로 렌더
        prst = self._read_prst(shape)
        elbow_points = self._bent_connector_points(prst, flip_h, flip_v)

        ln = spPr.find(qn("a:ln")) if spPr is not None else None
        tail_arrow = False
        head_arrow = False
        dash_style: str | None = None

        if ln is not None:
            tail_end = ln.find(qn("a:tailEnd"))
            if tail_end is not None and tail_end.get("type", "none") != "none":
                tail_arrow = True
            head_end = ln.find(qn("a:headEnd"))
            if head_end is not None and head_end.get("type", "none") != "none":
                head_arrow = True
            prst_dash = ln.find(qn("a:prstDash"))
            if prst_dash is not None:
                dash_val = prst_dash.get("val", "")
                # OOXML prstDash 값을 ("dash"|"dot") 으로 정규화
                if dash_val in (
                    "dash",
                    "lgDash",
                    "sysDash",
                    "dashDot",
                    "lgDashDot",
                    "lgDashDotDot",
                ):
                    dash_style = "dash"
                elif dash_val in ("dot", "sysDot"):
                    dash_style = "dot"

        width_px = self._emu_to_px_x(width_emu)
        height_px = self._emu_to_px_y(height_emu)

        if elbow_points is not None:
            # elbow 커넥터: 폴리라인 꼭짓점이 이미 flip 을 반영하므로 직선용 flip/방향
            # 로직을 건너뛴다. 화살표는 첫/끝 꼭짓점(start/end)에 그대로 매핑된다.
            return PptxShape(
                left_px=self._emu_to_px_x(left_emu),
                top_px=self._emu_to_px_y(top_emu),
                width_px=width_px,
                height_px=height_px,
                shape_type="line",
                border_color=border_color,
                border_width_pt=border_width,
                end_arrow=tail_arrow,
                start_arrow=head_arrow,
                dash_style=dash_style,
                elbow_points=elbow_points,
                rotation=self._read_rotation(shape),
            )

        # PPTX 커넥터의 실제 시작/끝 좌표 계산:
        # 기본: start=(left,top), end=(left+w,top+h)
        # flipH=1 → x좌표 반전: start.x=left+w, end.x=left
        # flipV=1 → y좌표 반전: start.y=top+h, end.y=top
        # headEnd 마커는 start에, tailEnd 마커는 end에 위치
        #
        # 디자인 스펙 모델은 항상 (left,top)→(left+w,top±h) 방향으로 그리므로
        # flip이 start/end를 뒤바꾸면 화살표도 swap 해야 함
        _SNAP = 12  # 수직/수평 판별 임계값 (px)
        is_horizontal = height_px == 0 or (width_px > 0 and abs(height_px) <= _SNAP)
        is_vertical = width_px == 0 or (abs(height_px) > 0 and width_px <= _SNAP)

        if is_vertical and not is_horizontal:
            # 수직선: flipH는 시각적 효과 없음, flipV만 방향 반전
            if flip_v:
                tail_arrow, head_arrow = head_arrow, tail_arrow
        elif is_horizontal and not is_vertical:
            # 수평선: flipV는 시각적 효과 없음, flipH만 방향 반전
            if flip_h:
                tail_arrow, head_arrow = head_arrow, tail_arrow
        else:
            # 대각선: flipH XOR flipV → 선 방향 반전 (↘ ↔ ↗)
            if flip_v != flip_h:
                height_px = -height_px
            # flipH가 활성이면 화살표 swap (↘→↙ 또는 ↗→↖)
            if flip_h:
                tail_arrow, head_arrow = head_arrow, tail_arrow

        return PptxShape(
            left_px=self._emu_to_px_x(left_emu),
            top_px=self._emu_to_px_y(top_emu),
            width_px=width_px,
            height_px=height_px,
            shape_type="line",
            border_color=border_color,
            border_width_pt=border_width,
            end_arrow=tail_arrow,
            start_arrow=head_arrow,
            dash_style=dash_style,
            rotation=self._read_rotation(shape),
        )

    @staticmethod
    def _bent_connector_points(
        prst: str | None, flip_h: bool, flip_v: bool
    ) -> list[list[float]] | None:
        """bentConnector prst → bbox 정규화 폴리라인 꼭짓점 (elbow). 아니면 None.

        OOXML 기본 기하(avLst 미지정=50%), unflipped 기준:
        - bentConnector2: L자, (0,0)→(1,0)→(1,1)  (수평 후 수직)
        - bentConnector3: Z자, (0,0)→(.5,0)→(.5,1)→(1,1)  (중간 50% 에서 꺾임)
        - bentConnector4/5 도 근사(중간 꺾임 2~3회)로 처리.
        flipH 는 x(fx→1-fx), flipV 는 y(fy→1-fy) 를 미러링한다.
        """
        if not prst or not prst.startswith("bentConnector"):
            return None
        base: list[tuple[float, float]]
        if prst == "bentConnector2":
            base = [(0.0, 0.0), (1.0, 0.0), (1.0, 1.0)]
        elif prst == "bentConnector3":
            base = [(0.0, 0.0), (0.5, 0.0), (0.5, 1.0), (1.0, 1.0)]
        elif prst in ("bentConnector4", "bentConnector5"):
            base = [
                (0.0, 0.0),
                (0.5, 0.0),
                (0.5, 0.5),
                (1.0, 0.5),
                (1.0, 1.0),
            ]
        else:
            return None
        pts = [
            [1.0 - fx if flip_h else fx, 1.0 - fy if flip_v else fy] for fx, fy in base
        ]
        return pts

    # ── 이미지 추출 ──

    def _extract_image(self, shape) -> PptxImage:
        try:
            blob = shape.image.blob
        except Exception:
            logger.debug(
                "이미지 blob 추출 실패 (name=%s)",
                getattr(shape, "name", "?"),
                exc_info=True,
            )
            blob = b""

        svg_bytes = _extract_svg_blip_bytes(shape)
        if blob and svg_bytes:
            blob = _recolor_svg_theme_fallback(
                blob,
                svg_bytes,
                getattr(self, "_theme_color_map", {}),
            )

        width_px = self._emu_to_px_x(shape.width)
        height_px = self._emu_to_px_y(shape.height)

        # prstGeom 으로 이미지 클리핑 모양 추출 → corner_radius 로 표현
        # ellipse=원형(반경=단변/2), roundRect=둥근모서리(고정 반경)
        corner_radius = None
        try:
            spPr = shape._element.find(qn("p:spPr"))
            geom = spPr.find(qn("a:prstGeom")) if spPr is not None else None
            prst = geom.get("prst") if geom is not None else None
            if prst == "ellipse":
                corner_radius = min(width_px, height_px) / 2
            elif prst == "roundRect":
                corner_radius = min(width_px, height_px) * 0.08
        except Exception:
            logger.debug("이미지 prstGeom 추출 실패", exc_info=True)

        return PptxImage(
            left_px=self._emu_to_px_x(shape.left),
            top_px=self._emu_to_px_y(shape.top),
            width_px=width_px,
            height_px=height_px,
            image_bytes=blob,
            corner_radius_px=corner_radius,
        )
