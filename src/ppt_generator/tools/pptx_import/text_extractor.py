"""텍스트/런/불릿 추출 모듈.

SlideReader에서 사용하는 텍스트 관련 추출 로직을 담당한다.
paragraph, run, 불릿 레벨, 정렬, 줄간격, 수직 정렬 등의 추출 기능을 제공한다.
"""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from pptx.oxml.ns import qn

from ppt_generator.interfaces.constants import (
    PPTX_BULLET_MARGIN_EMU_L1,
)
from ppt_generator.interfaces.schemas import (
    PptxParagraph,
    PptxTextRun,
)
from ppt_generator.tools.pptx_import.ooxml_utils import (
    ALIGN_REVERSE_MAP,
    MONOSPACE_FONTS,
)
from ppt_generator.tools.pptx_import.theme_resolver import (
    PH_TYPE_TO_TXSTYLE,
    DefaultRunProps,
    extract_color_from_rpr,
    extract_props_from_rpr,
)

if TYPE_CHECKING:
    pass

logger = logging.getLogger(__name__)


class TextExtractorMixin:
    """텍스트 추출 기능을 제공하는 mixin 클래스.

    SlideReader에서 상속하여 사용한다.
    _theme_color_map, _master_tx_styles, _layout_def_rpr 속성이 필요하다.
    """

    _theme_color_map: dict[str, str]
    _master_tx_styles: dict[str, dict[int, DefaultRunProps]]
    _layout_def_rpr: dict[int, dict[int, DefaultRunProps]]

    def _extract_paragraphs(
        self,
        text_frame,
        placeholder_type: int | None = None,
        placeholder_idx: int | None = None,
    ) -> list[PptxParagraph]:
        paragraphs: list[PptxParagraph] = []
        for para in text_frame.paragraphs:
            bullet_level = self._extract_bullet_level(para)
            runs = self._extract_runs(
                para,
                placeholder_type=placeholder_type,
                placeholder_idx=placeholder_idx,
                bullet_level=max(bullet_level, 0),
            )
            if not runs:
                continue
            alignment = self._extract_alignment(para)
            paragraphs.append(
                PptxParagraph(
                    runs=runs,
                    bullet_level=bullet_level,
                    alignment=alignment,
                )
            )
        return paragraphs

    def _resolve_inherited_props(
        self,
        paragraph,
        placeholder_type: int | None,
        placeholder_idx: int | None,
        bullet_level: int,
    ) -> DefaultRunProps:
        """OOXML 상속 체인에서 paragraph → layout → master 순으로 기본 서식을 resolve."""
        font_size: int | None = None
        color: str | None = None
        bold: bool | None = None
        font_name: str | None = None

        # 1) paragraph a:pPr > a:defRPr
        pPr = paragraph._p.find(qn("a:pPr"))
        if pPr is not None:
            def_rPr = pPr.find(qn("a:defRPr"))
            para_props = extract_props_from_rpr(def_rPr, self._theme_color_map)
            if para_props.font_size_pt is not None:
                font_size = para_props.font_size_pt
            if para_props.color is not None:
                color = para_props.color
            if para_props.bold is not None:
                bold = para_props.bold
            if para_props.font_name is not None:
                font_name = para_props.font_name

        # 2) layout placeholder lstStyle > lvlNpPr > defRPr
        if placeholder_idx is not None and placeholder_idx in self._layout_def_rpr:
            layout_levels = self._layout_def_rpr[placeholder_idx]
            layout_props = layout_levels.get(bullet_level)
            if layout_props is not None:
                if font_size is None and layout_props.font_size_pt is not None:
                    font_size = layout_props.font_size_pt
                if color is None and layout_props.color is not None:
                    color = layout_props.color
                if bold is None and layout_props.bold is not None:
                    bold = layout_props.bold
                if font_name is None and layout_props.font_name is not None:
                    font_name = layout_props.font_name

        # 3) master p:txStyles > titleStyle/bodyStyle/otherStyle > lvlNpPr > defRPr
        style_name = (
            PH_TYPE_TO_TXSTYLE.get(placeholder_type, "otherStyle")
            if placeholder_type is not None
            else "otherStyle"
        )
        master_levels = self._master_tx_styles.get(style_name, {})
        master_props = master_levels.get(bullet_level)
        if master_props is not None:
            if font_size is None and master_props.font_size_pt is not None:
                font_size = master_props.font_size_pt
            if color is None and master_props.color is not None:
                color = master_props.color
            if bold is None and master_props.bold is not None:
                bold = master_props.bold
            if font_name is None and master_props.font_name is not None:
                font_name = master_props.font_name

        return DefaultRunProps(
            font_size_pt=font_size, color=color, bold=bold, font_name=font_name
        )

    def _extract_runs(
        self,
        paragraph,
        placeholder_type: int | None = None,
        placeholder_idx: int | None = None,
        bullet_level: int = 0,
    ) -> list[PptxTextRun]:
        inherited = self._resolve_inherited_props(
            paragraph,
            placeholder_type,
            placeholder_idx,
            bullet_level,
        )

        # paragraph.runs 는 <a:r> 만 반환하고 <a:br>(문단 내 소프트 줄바꿈)과
        # <a:fld>(슬라이드 번호·날짜 등 필드)를 건너뛴다. 그대로 이어붙이면
        # "Complex"+"re-architecting" 처럼 단어가 붙거나 필드 텍스트가 사라지므로,
        # XML 자식을 순서대로 순회해 <a:br> 위치에 개행을, <a:fld> 를 run 으로 보존한다.
        # <a:r> 요소는 paragraph.runs 와 문서 순서가 동일하므로 인덱스로 매칭한다.
        _R = qn("a:r")
        _BR = qn("a:br")
        _FLD = qn("a:fld")
        pptx_runs = list(paragraph.runs)
        run_idx = 0
        pending_break = False

        runs: list[PptxTextRun] = []
        for child in paragraph._p:
            if child.tag == _BR:
                pending_break = True
                continue
            if child.tag == _FLD:
                # 필드(slidenum/datetime 등): <a:t> 에 마지막 렌더값이 들어있다.
                fld_run = self._field_to_run(child, inherited, pending_break)
                pending_break = False
                if fld_run is not None:
                    runs.append(fld_run)
                continue
            if child.tag != _R:
                continue
            run = pptx_runs[run_idx] if run_idx < len(pptx_runs) else None
            run_idx += 1
            if run is None or not run.text:
                pending_break = False
                continue
            font = run.font

            font_size: int | None = None
            if font.size is not None:
                font_size = round(font.size.pt)
            if font_size is None:
                font_size = inherited.font_size_pt

            color: str | None = None
            try:
                if font.color and font.color.rgb:
                    color = f"#{font.color.rgb}"
            except (AttributeError, TypeError):
                logger.debug("폰트 색상 추출 실패", exc_info=True)
            if color is None:
                try:
                    rPr = run._r.find(qn("a:rPr"))
                    if rPr is not None:
                        color = extract_color_from_rpr(rPr, self._theme_color_map)
                except Exception:
                    logger.debug("rPr 색상 추출 실패", exc_info=True)
            if color is None:
                color = inherited.color

            run_bold = font.bold
            if run_bold is None:
                run_bold = inherited.bold if inherited.bold is not None else False
            else:
                run_bold = bool(run_bold)

            font_family: str | None = None
            if font.name and font.name.lower() in MONOSPACE_FONTS:
                font_family = "monospace"

            # 원본 폰트명 보존 (상속 우선순위):
            # run 명시 latin → layout/master placeholder 상속(inherited.font_name)
            # → 테마 폰트 폴백. 상속 단계를 건너뛰고 곧장 테마 major 로 가면
            # layout 이 지정한 실제 폰트(예: title 의 "Amazon Ember Display")를 놓치고
            # theme major("...Heavy")로 잘못 폴백해 글자 폭이 달라진다 (import/0003).
            font_name = font.name or inherited.font_name
            if not font_name:
                theme_fonts = getattr(self, "_theme_fonts", {})
                if placeholder_type in (1, 3, 15):
                    font_name = theme_fonts.get("major") or theme_fonts.get("minor")
                else:
                    font_name = theme_fonts.get("minor") or theme_fonts.get("major")

            href: str | None = None
            try:
                if run.hyperlink and run.hyperlink.address:
                    href = run.hyperlink.address
            except Exception:
                logger.debug("하이퍼링크 추출 실패", exc_info=True)

            # 직전에 <a:br> 이 있었다면 이 run 앞에 개행을 붙여 소프트 줄바꿈 보존.
            run_text = run.text
            if pending_break:
                run_text = "\n" + run_text
                pending_break = False

            runs.append(
                PptxTextRun(
                    text=run_text,
                    font_size_pt=font_size,
                    color=color,
                    bold=run_bold,
                    italic=bool(font.italic),
                    font_family=font_family,
                    href=href,
                    font_name=font_name,
                )
            )
        return runs

    def _field_to_run(
        self, fld, inherited: DefaultRunProps, pending_break: bool
    ) -> PptxTextRun | None:
        """<a:fld>(slidenum/datetime 등 필드)를 PptxTextRun 으로 변환.

        필드는 <a:rPr> + <a:t>(마지막 렌더값) 구조로 <a:r> 과 사실상 동일하다.
        """
        t_el = fld.find(qn("a:t"))
        text = t_el.text if t_el is not None else None
        if not text:
            return None
        if pending_break:
            text = "\n" + text

        rPr = fld.find(qn("a:rPr"))
        props = extract_props_from_rpr(rPr, self._theme_color_map)
        font_size = props.font_size_pt or inherited.font_size_pt
        color = props.color or inherited.color
        bold = props.bold if props.bold is not None else (inherited.bold or False)

        font_name: str | None = None
        latin = rPr.find(qn("a:latin")) if rPr is not None else None
        if latin is not None:
            font_name = latin.get("typeface")
        if not font_name:
            theme_fonts = getattr(self, "_theme_fonts", {})
            font_name = theme_fonts.get("minor") or theme_fonts.get("major")

        font_family: str | None = None
        if font_name and font_name.lower() in MONOSPACE_FONTS:
            font_family = "monospace"

        return PptxTextRun(
            text=text,
            font_size_pt=font_size,
            color=color,
            bold=bool(bold),
            italic=False,
            font_family=font_family,
            font_name=font_name,
        )

    @staticmethod
    def _extract_bullet_level(paragraph) -> int:
        pPr = paragraph._p.find(qn("a:pPr"))
        if pPr is None:
            return -1

        has_bullet = (
            pPr.find(qn("a:buChar")) is not None
            or pPr.find(qn("a:buAutoNum")) is not None
        )
        if pPr.find(qn("a:buNone")) is not None:
            has_bullet = False

        if not has_bullet:
            lvl = pPr.get("lvl")
            if lvl is not None:
                lvl_int = int(lvl)
                if lvl_int > 0:
                    return lvl_int
            return -1

        lvl = pPr.get("lvl")
        if lvl is not None:
            return int(lvl)

        marL_str = pPr.get("marL")
        if marL_str is not None:
            marL = int(marL_str)
            if marL >= PPTX_BULLET_MARGIN_EMU_L1:
                return 1
        return 0

    @staticmethod
    def _extract_alignment(paragraph) -> str | None:
        # paragraph.alignment 가 None 이면 PPTX 의 master/layout 기본값에 따라 렌더되지만,
        # 우리 spec 은 master inheritance 를 추적하지 않으므로 동일 spec 이 PPTX/HTML 에서
        # 다르게 렌더된다 (PPTX 는 master 의 algn=ctr, HTML 은 CSS 기본 left). 좌측으로
        # 정규화해 두 렌더러가 일관되게 좌측 정렬되도록 한다.
        if paragraph.alignment is None:
            return "left"
        return ALIGN_REVERSE_MAP.get(paragraph.alignment) or "left"

    def _extract_line_spacing(
        self,
        text_frame,
        placeholder_type: int | None = None,
        placeholder_idx: int | None = None,
    ) -> float | None:
        """텍스트 프레임의 줄간격(pt)을 추출. 첫 번째 paragraph 기준.

        1) 문단에 직접 지정된 줄간격: 절대(pt)면 그대로, 배수(예 0.9)면 폰트 크기와 곱해 pt 환산.
        2) 직접 지정이 없으면(placeholder 상속) layout lstStyle → master txStyle 의
           lnSpc 를 해석한다. spcPct(배수)는 상속 폰트 크기와 곱해 pt 로 환산한다.
           원본 마스터가 90% 줄간격을 지정했는데 이를 놓치면 렌더러 기본 1.5 가 적용돼
           텍스트가 세로로 퍼져 박스를 넘친다 (import/0003).
        """
        for para in text_frame.paragraphs:
            spacing = para.line_spacing
            if spacing is not None:
                # 절대 pt (Length) — .pt 접근 가능
                try:
                    return float(spacing.pt)
                except (AttributeError, TypeError):
                    pass
                # 배수 (float) — 상속 폰트 크기와 곱해 pt 로 환산
                if isinstance(spacing, (int, float)):
                    font_pt = self._first_run_font_pt(
                        para, placeholder_type, placeholder_idx
                    )
                    if font_pt:
                        return float(spacing) * font_pt
            break  # 첫 문단만 기준 (기존 동작 유지)

        # placeholder 상속 lnSpc 해석
        return self._resolve_inherited_line_spacing(placeholder_type, placeholder_idx)

    def _first_run_font_pt(
        self,
        paragraph,
        placeholder_type: int | None,
        placeholder_idx: int | None,
    ) -> int | None:
        """문단 첫 run 의 유효 폰트 크기(pt). 직접 지정 → 상속 순으로 해석."""
        for run in paragraph.runs:
            sz = run.font.size
            if sz is not None:
                try:
                    return round(sz.pt)
                except (AttributeError, TypeError):
                    pass
            break
        inherited = self._resolve_inherited_props(
            paragraph, placeholder_type, placeholder_idx, 0
        )
        return inherited.font_size_pt

    def _resolve_inherited_line_spacing(
        self,
        placeholder_type: int | None,
        placeholder_idx: int | None,
    ) -> float | None:
        """layout lstStyle → master txStyle 상속 체인에서 레벨 0 lnSpc 를 pt 로 해석."""
        for props in (
            self._layout_def_rpr.get(placeholder_idx, {}).get(0)
            if placeholder_idx is not None
            else None,
            self._master_tx_styles.get(
                PH_TYPE_TO_TXSTYLE.get(placeholder_type, "otherStyle")
                if placeholder_type is not None
                else "otherStyle",
                {},
            ).get(0),
        ):
            if props is None:
                continue
            if props.line_spacing_pt is not None:
                return props.line_spacing_pt
            if props.line_spacing_pct is not None and props.font_size_pt is not None:
                return props.line_spacing_pct * props.font_size_pt
        return None

    def _extract_vertical_alignment(
        self, text_frame, placeholder_idx: int | None = None
    ) -> str:
        anchor_map = {"t": "top", "ctr": "middle", "b": "bottom"}
        try:
            bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
            if bodyPr is not None:
                anchor = bodyPr.get("anchor")
                if anchor in anchor_map:
                    return anchor_map[anchor]
            # shape 자체에 anchor 가 없으면 placeholder 상속(layout→master)을 따른다.
            # PowerPoint 는 예컨대 master TITLE placeholder 의 anchor="ctr" 를
            # 상속하므로, 이를 무시하면 제목이 top 으로 렌더돼 세로 위치가 어긋난다.
            if placeholder_idx is not None:
                inherited = getattr(self, "_ph_anchor_by_idx", {}).get(placeholder_idx)
                if inherited in anchor_map:
                    return anchor_map[inherited]
        except Exception:
            logger.debug("vertical_alignment 추출 실패", exc_info=True)
        return "top"

    @staticmethod
    def _extract_autofit(text_frame) -> tuple[str, float | None]:
        """<a:bodyPr> 의 autofit 자식에서 (모드, fontScale) 을 추출.

        - <a:noAutofit/>  → ("none", None)   : 축소 없이 넘침 허용
        - <a:spAutoFit/>  → ("resize", None)  : 박스가 텍스트에 맞춰 커짐
        - <a:normAutofit fontScale="90000"/> → ("none", 0.9)
              PPT 가 저장한 실제 축소율을 그대로 적용 (fontScale 없으면 100%).
              PowerPoint 는 normAutofit 을 재계산하지 않고 저장된 스케일을 쓰므로,
              우리도 렌더러 재계산 대신 이 값을 적용해야 원본과 일치한다.
        - 없음 → ("none", None) : OOXML 명세상 autofit 자식이 없으면 기본은
              noAutofit(축소 없이 넘침 허용)이다. 임포트는 원본 레이아웃이 이미
              확정된 상태이므로 렌더러가 폰트를 재축소하면 원본보다 작아진다
              (import/0001 "임포트 시 autofit 비활성화" 결정). LLM 생성 슬라이드의
              기본값("shrink")과 반대라는 점에 주의 — 이 함수는 임포트 전용이다.
        """
        try:
            bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
            if bodyPr is not None:
                if bodyPr.find(qn("a:noAutofit")) is not None:
                    return "none", None
                if bodyPr.find(qn("a:spAutoFit")) is not None:
                    return "resize", None
                norm = bodyPr.find(qn("a:normAutofit"))
                if norm is not None:
                    fs = norm.get("fontScale")
                    scale = int(fs) / 100000 if fs else None
                    # normAutofit 은 PPT 저장 스케일을 그대로 쓰므로 재계산 안 함
                    return "none", scale
        except Exception:
            logger.debug("autofit 추출 실패", exc_info=True)
        return "none", None
