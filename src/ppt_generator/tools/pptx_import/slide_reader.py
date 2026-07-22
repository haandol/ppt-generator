"""python-pptx 슬라이드 객체에서 PptxSlideSpec을 추출하는 리더 모듈.

SlideBuilder의 역변환: python-pptx 오브젝트 → PptxSlideSpec 데이터클래스.

상세 로직은 하위 모듈로 분리:
- theme_resolver: 테마/색상 해석
- ooxml_utils: OOXML 상수 매핑, custGeom → SVG 변환
- text_extractor: 텍스트/런/불릿 추출
- style_extractor: fill/line/dash 스타일 추출
- shape_extractors: 개별 도형 유형별 추출
"""

from __future__ import annotations

import logging
from dataclasses import replace
from typing import TYPE_CHECKING

from pptx.oxml.ns import qn

from ppt_generator.interfaces.constants import (
    IMPORT_EMU_TO_PX,
    PX_TO_EMU,
    SLIDES_HEIGHT_PX,
    SLIDES_WIDTH_PX,
)
from ppt_generator.interfaces.schemas import (
    PptxImage,
    PptxShape,
    PptxSlideSpec,
    PptxTextBox,
)
from ppt_generator.tools.pptx_import.ooxml_utils import CLOSING_KEYWORDS
from ppt_generator.tools.pptx_import.shape_extractors import ShapeExtractorMixin
from ppt_generator.tools.pptx_import.style_extractor import StyleExtractorMixin
from ppt_generator.tools.pptx_import.text_extractor import TextExtractorMixin
from ppt_generator.tools.pptx_import.theme_resolver import (
    DefaultRunProps,
    extract_line_spacing_from_ppr,
    extract_master_tx_styles,
    extract_props_from_rpr,
    extract_theme_color_map,
    extract_theme_color_map_for_master,
    extract_theme_fonts_for_master,
    resolve_scheme_color,
)

# 하위 호환: 외부에서 기존 private 이름으로 import하는 코드 지원
_DefaultRunProps = DefaultRunProps
_extract_theme_color_map = extract_theme_color_map

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide

logger = logging.getLogger(__name__)


class SlideReader(ShapeExtractorMixin, TextExtractorMixin, StyleExtractorMixin):
    """python-pptx 슬라이드에서 PptxSlideSpec을 추출하는 리더."""

    def __init__(
        self,
        scale_x: float = 1.0,
        scale_y: float = 1.0,
        presentation: Presentation | None = None,
    ) -> None:
        self._scale_x = scale_x
        self._scale_y = scale_y
        self._theme_color_map: dict[str, str] = {}
        self._master_tx_styles: dict[str, dict[int, DefaultRunProps]] = {}
        self._presentation_width: int = 0
        self._presentation_height: int = 0
        # 멀티 마스터 지원: 마스터별 색상맵/서식을 캐시하고, 슬라이드마다 소속
        # 마스터에 맞는 것을 활성화한다. (첫 마스터 값은 기본 폴백으로 유지)
        self._theme_map_by_master: dict[int, dict[str, str]] = {}
        self._tx_styles_by_master: dict[int, dict[str, dict[int, DefaultRunProps]]] = {}
        self._theme_fonts_by_master: dict[int, dict[str, str]] = {}
        # 현재 활성 슬라이드의 테마 폰트 (major/minor) — run 에 명시 폰트가 없을 때 폴백
        self._theme_fonts: dict[str, str] = {}
        if presentation is not None:
            self._presentation_width = presentation.slide_width or 0
            self._presentation_height = presentation.slide_height or 0
            self._theme_color_map = extract_theme_color_map(presentation)
            self._master_tx_styles = extract_master_tx_styles(presentation)
            for master in presentation.slide_masters:
                key = id(master._element)
                self._theme_map_by_master[key] = extract_theme_color_map_for_master(
                    master
                )
                self._tx_styles_by_master[key] = extract_master_tx_styles(
                    presentation, master
                )
                self._theme_fonts_by_master[key] = extract_theme_fonts_for_master(
                    master
                )
        self._layout_def_rpr: dict[int, dict[int, DefaultRunProps]] = {}
        self._default_bg_color: str | None = None

    def set_default_bg_color(self, color: str | None) -> None:
        """기본 배경색을 설정 (최종 폴백으로 사용)."""
        self._default_bg_color = color

    @staticmethod
    def compute_scale(presentation: Presentation) -> tuple[float, float]:
        """프레젠테이션 슬라이드 크기 → 1280×720 캔버스 스케일 팩터 계산."""
        src_width_px = presentation.slide_width * IMPORT_EMU_TO_PX
        src_height_px = presentation.slide_height * IMPORT_EMU_TO_PX
        scale_x = SLIDES_WIDTH_PX / src_width_px if src_width_px else 1.0
        scale_y = SLIDES_HEIGHT_PX / src_height_px if src_height_px else 1.0
        return scale_x, scale_y

    # ── 좌표 변환 ──

    def _emu_to_px_x(self, emu: int) -> float:
        return round(emu * IMPORT_EMU_TO_PX * self._scale_x, 1)

    def _emu_to_px_y(self, emu: int) -> float:
        return round(emu * IMPORT_EMU_TO_PX * self._scale_y, 1)

    def _emu_to_px_padding(self, emu: int | None) -> float | None:
        if emu is None:
            return None
        return round(emu / PX_TO_EMU, 1)

    # ── 슬라이드 읽기 ──

    def read_slide(
        self,
        slide: Slide,
        slide_index: int,
        total_slides: int,
    ) -> PptxSlideSpec:
        """단일 슬라이드 → PptxSlideSpec 변환."""
        self._activate_master_theme(slide)
        self._cache_layout_def_rpr(slide)
        background_color = self._extract_background_color(slide)
        background_image_bytes = self._extract_background_image_bytes(slide)
        textboxes: list[PptxTextBox] = []
        shapes: list[PptxShape] = []
        images: list[PptxImage] = []
        warnings: list[str] = []

        z_counter = 0

        # 1) 레이아웃/마스터의 정적 요소(로고, "Thank you!", 저작권 등 비-placeholder
        #    텍스트박스·그림·도형)를 먼저 상속한다. placeholder 는 슬라이드가 오버라이드
        #    하므로 제외하고, full-bleed 배경 그림도 배경으로 이미 처리하므로 제외한다.
        # 슬라이드 자체 텍스트를 미리 등록해, 레이아웃·마스터에서 동일 텍스트(저작권 등)를
        # 중복 상속하지 않도록 한다.
        seen_static_text: set = set()
        try:
            for sh in slide.shapes:
                if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                    seen_static_text.add("".join(sh.text_frame.text.split()))
        except Exception:
            logger.debug("슬라이드 텍스트 사전 스캔 실패", exc_info=True)

        for source in (
            getattr(slide, "slide_layout", None),
            getattr(getattr(slide, "slide_layout", None), "slide_master", None),
        ):
            if source is None:
                continue
            for shape in self._static_inherited_shapes(source):
                # 레이아웃·마스터 중복 방지: 동일 텍스트 요소는 한 번만 상속
                if (
                    getattr(shape, "has_text_frame", False)
                    and shape.text_frame.text.strip()
                ):
                    key = "".join(shape.text_frame.text.split())
                    if key in seen_static_text:
                        continue
                    seen_static_text.add(key)
                else:
                    # 그림 등: 위치+크기로 중복 판정 (같은 로고가 레이아웃·마스터에 중복)
                    pos_key = (
                        round((shape.left or 0) / 9525),
                        round((shape.top or 0) / 9525),
                        round((shape.width or 0) / 9525),
                        round((shape.height or 0) / 9525),
                    )
                    if pos_key in seen_static_text:
                        continue
                    seen_static_text.add(pos_key)
                prev_counts = (len(textboxes), len(shapes), len(images))
                try:
                    self._extract_shape(shape, textboxes, shapes, images, warnings)
                except Exception:
                    logger.debug(
                        "상속 요소 추출 실패 (name=%s)",
                        getattr(shape, "name", "?"),
                        exc_info=True,
                    )
                z_counter = self._assign_z(
                    textboxes, shapes, images, prev_counts, z_counter
                )

        # 2) 슬라이드 자체 요소 (상속 요소 위에 렌더됨)
        for shape in slide.shapes:
            prev_counts = (len(textboxes), len(shapes), len(images))
            try:
                self._extract_shape(
                    shape,
                    textboxes,
                    shapes,
                    images,
                    warnings,
                )
            except Exception:
                logger.warning(
                    "슬라이드 %d: shape 추출 실패 (name=%s, type=%s)",
                    slide_index + 1,
                    getattr(shape, "name", "?"),
                    getattr(shape, "shape_type", "?"),
                    exc_info=True,
                )
            # 새로 추가된 요소에 z_index 할당 (그룹 내부는 draw-order 존중)
            z_counter = self._assign_z(
                textboxes, shapes, images, prev_counts, z_counter
            )

        speaker_notes = self._extract_speaker_notes(slide)
        slide_type = self._infer_slide_type(
            slide_index,
            total_slides,
            textboxes,
            shapes,
        )

        for w in warnings:
            logger.warning("슬라이드 %d: %s", slide_index + 1, w)

        return PptxSlideSpec(
            background_color=background_color,
            background_image_bytes=background_image_bytes,
            textboxes=textboxes,
            shapes=shapes,
            images=images,
            speaker_notes=speaker_notes,
            slide_type=slide_type,
        )

    def _assign_z(
        self,
        textboxes: list[PptxTextBox],
        shapes: list[PptxShape],
        images: list[PptxImage],
        prev_counts: tuple[int, int, int],
        z_counter: int,
    ) -> int:
        """한 top-level shape 가 새로 추가한 요소들에 전역 z_index 를 부여한다.

        그룹은 자식 요소에 그룹 내 상대 순서(draw-order)를 z_index 로 남겨 둔다
        (compound_extractors). 이를 정렬 키로 삼아 그리기 순서를 보존한 채 전역 z 를
        매긴다. 상대 순서가 없는(단일) 요소는 리스트 종류 순서로 안정 정렬된다.
        z_index 를 리스트 종류별로만 매기면 "박스 뒤 + 텍스트 앞" 그룹에서 순서가 뒤집혀
        흰 박스가 텍스트를 덮는다 (import/0003).
        """
        new_items: list[tuple[int, int, list, int]] = []
        for kind, (lst, prev_len) in enumerate(
            (
                (textboxes, prev_counts[0]),
                (shapes, prev_counts[1]),
                (images, prev_counts[2]),
            )
        ):
            for idx in range(prev_len, len(lst)):
                rel = lst[idx].z_index  # 그룹 내 상대 순서 (없으면 None)
                new_items.append((rel if rel is not None else 0, kind, lst, idx))

        # (그룹 상대순서, 리스트종류) 로 안정 정렬 → 그리기 순서 보존
        new_items.sort(key=lambda t: (t[0], t[1]))
        for _rel, _kind, lst, idx in new_items:
            lst[idx] = replace(lst[idx], z_index=z_counter)
            z_counter += 1
        return z_counter

    def _static_inherited_shapes(self, source) -> list:
        """레이아웃/마스터에서 슬라이드로 상속할 정적 요소를 선별.

        - placeholder 는 제외 (슬라이드가 값을 오버라이드)
        - full-bleed 배경 그림은 제외 (배경으로 이미 처리)
        - 텍스트 없는 빈 도형/장식은 제외 (노이즈 방지)
        - 로고 그림, 정적 텍스트박스("Thank you!", 저작권 등)만 포함
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        result = []
        slide_w = self._presentation_width
        slide_h = self._presentation_height
        try:
            for shape in source.shapes:
                if getattr(shape, "is_placeholder", False):
                    continue
                st = getattr(shape, "shape_type", None)
                if st == MSO_SHAPE_TYPE.PICTURE:
                    # full-bleed 배경 그림 제외
                    w = shape.width or 0
                    h = shape.height or 0
                    if (
                        slide_w
                        and slide_h
                        and w >= slide_w * 0.95
                        and h >= slide_h * 0.95
                    ):
                        continue
                    result.append(shape)
                    continue
                # 텍스트가 있는 도형/텍스트박스만 (빈 장식 도형 제외)
                if getattr(shape, "has_text_frame", False):
                    if shape.text_frame.text.strip():
                        result.append(shape)
        except Exception:
            logger.debug("상속 요소 선별 실패", exc_info=True)
        return result

    # ── 마스터 테마 활성화 ──

    def _activate_master_theme(self, slide: Slide) -> None:
        """슬라이드가 소속된 마스터의 색상맵/서식을 활성화한다.

        멀티 마스터(다크/라이트 혼재) 프레젠테이션에서 각 슬라이드가 올바른
        테마 색상으로 배경/텍스트를 해석하도록 한다.
        """
        try:
            master = slide.slide_layout.slide_master
            key = id(master._element)
            theme_map = self._theme_map_by_master.get(key)
            if theme_map:
                self._theme_color_map = theme_map
            tx_styles = self._tx_styles_by_master.get(key)
            if tx_styles is not None:
                self._master_tx_styles = tx_styles
            self._theme_fonts = self._theme_fonts_by_master.get(key, {})
        except Exception:
            logger.debug("마스터 테마 활성화 실패", exc_info=True)

    # ── 레이아웃 캐시 ──

    def _cache_layout_def_rpr(self, slide: Slide) -> None:
        """현재 슬라이드의 layout placeholder별 lstStyle > defRPr을 캐시."""
        self._layout_def_rpr = {}
        try:
            layout = slide.slide_layout
            for ph in layout.placeholders:
                ph_idx = ph.placeholder_format.idx
                ph_el = ph._element
                txBody = ph_el.find(qn("p:txBody"))
                if txBody is None:
                    continue
                lstStyle = txBody.find(qn("a:lstStyle"))
                if lstStyle is None:
                    continue
                level_map: dict[int, DefaultRunProps] = {}
                for lvl_idx in range(9):
                    lvl_pPr = lstStyle.find(qn(f"a:lvl{lvl_idx + 1}pPr"))
                    if lvl_pPr is None:
                        continue
                    def_rPr = lvl_pPr.find(qn("a:defRPr"))
                    if def_rPr is not None:
                        props = extract_props_from_rpr(def_rPr, self._theme_color_map)
                        pct, pts = extract_line_spacing_from_ppr(lvl_pPr)
                        props.line_spacing_pct = pct
                        props.line_spacing_pt = pts
                        level_map[lvl_idx] = props
                if level_map:
                    self._layout_def_rpr[ph_idx] = level_map
        except Exception:
            logger.debug("레이아웃 defRPr 캐시 실패", exc_info=True)

    # ── 배경 추출 ──

    def _extract_background_color(self, slide: Slide) -> str | None:
        color = self._try_extract_bg_fill(slide.background)
        if color:
            return color

        try:
            color = self._try_extract_bg_fill(slide.slide_layout.background)
            if color:
                return color
        except Exception:
            logger.debug("레이아웃 배경색 추출 실패", exc_info=True)

        try:
            color = self._try_extract_bg_fill(
                slide.slide_layout.slide_master.background
            )
            if color:
                return color
        except Exception:
            logger.debug("마스터 배경색 추출 실패", exc_info=True)

        return self._default_bg_color or self._theme_color_map.get("bg1")

    def _extract_background_image_bytes(self, slide: Slide) -> bytes:
        """슬라이드 배경 이미지 바이트를 추출.

        slide → layout → master 순서로 탐색하며, 각 소스에서 두 가지 형태를 본다:
        1) <p:bg> 의 blipFill (명시적 배경 이미지)
        2) full-bleed 그림 도형 (레이아웃/마스터에 배경처럼 깔린 <pic>)

        많은 템플릿이 그라데이션/브랜드 배경을 <p:bg> 가 아니라 슬라이드 전체를 덮는
        그림 도형으로 넣기 때문에 후자도 배경으로 취급해야 원본과 일치한다.
        """
        for source in (
            slide,
            getattr(slide, "slide_layout", None),
            getattr(getattr(slide, "slide_layout", None), "slide_master", None),
        ):
            if source is None:
                continue
            # 1) <p:bg> blipFill
            blob = self._bg_blip_bytes(source)
            if blob:
                return blob
            # 2) full-bleed 그림 도형 (slide 자체는 사용자 콘텐츠이므로 제외,
            #    레이아웃/마스터에서만 배경 그림으로 간주)
            if source is not slide:
                blob = self._full_bleed_picture_bytes(source)
                if blob:
                    return blob
        return b""

    def _bg_blip_bytes(self, source) -> bytes:
        """source.background 의 <p:bg> blipFill 이미지 바이트."""
        try:
            bg = source.background
            bg_el = bg._element if hasattr(bg, "_element") else bg
            p_bg = bg_el.find(f".//{qn('p:bg')}")
            if p_bg is None:
                return b""
            blip = p_bg.find(f".//{qn('a:blip')}")
            if blip is None:
                return b""
            embed_id = blip.get(qn("r:embed"))
            if not embed_id:
                return b""
            part = getattr(source, "part", None)
            if part is None:
                return b""
            return part.related_part(embed_id).blob
        except Exception:
            logger.debug("배경 blipFill 추출 실패", exc_info=True)
            return b""

    def _full_bleed_picture_bytes(self, source) -> bytes:
        """레이아웃/마스터에서 슬라이드 전체를 덮는 그림 도형의 이미지 바이트.

        좌상단 근처에서 시작해 캔버스의 대부분(≥95%)을 덮는 그림을 배경으로 본다.
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            slide_w = self._presentation_width
            slide_h = self._presentation_height
            if not slide_w or not slide_h:
                return b""
            for shape in source.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                left = shape.left or 0
                top = shape.top or 0
                width = shape.width or 0
                height = shape.height or 0
                covers = (
                    width >= slide_w * 0.95
                    and height >= slide_h * 0.95
                    and left <= slide_w * 0.05
                    and top <= slide_h * 0.05
                )
                if not covers:
                    continue
                try:
                    return shape.image.blob
                except Exception:
                    logger.debug("full-bleed 그림 blob 추출 실패", exc_info=True)
            return b""
        except Exception:
            logger.debug("full-bleed 그림 배경 추출 실패", exc_info=True)
            return b""

    def _try_extract_bg_fill(self, bg) -> str | None:
        """background 객체에서 solid fill 색상을 XML 직접 파싱으로 추출."""
        try:
            bg_el = bg._element if hasattr(bg, "_element") else bg
            p_bg = bg_el.find(qn("p:bg")) if bg_el.tag != qn("p:bg") else bg_el
            if p_bg is None:
                return None
            if p_bg.find(f".//{qn('a:noFill')}") is not None:
                return None
            for srgb in p_bg.iter(qn("a:srgbClr")):
                val = srgb.get("val")
                if val:
                    return f"#{val}"
            for scheme_el in p_bg.iter(qn("a:schemeClr")):
                return resolve_scheme_color(scheme_el, self._theme_color_map)
        except Exception:
            logger.debug("배경 fill 추출 실패", exc_info=True)
        return None

    # ── 발표자 노트 ──

    @staticmethod
    def _extract_speaker_notes(slide: Slide) -> str:
        try:
            if slide.has_notes_slide:
                return slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            logger.debug("발표자 노트 추출 실패", exc_info=True)
        return ""

    # ── slide_type 추론 ──

    @staticmethod
    def _infer_slide_type(
        slide_index: int,
        total_slides: int,
        textboxes: list[PptxTextBox],
        shapes: list[PptxShape],
    ) -> str:
        all_text = ""
        max_font = 0
        element_count = len(textboxes) + len(shapes)

        for tb in textboxes:
            for p in tb.paragraphs:
                for r in p.runs:
                    all_text += r.text
                    if r.font_size_pt and r.font_size_pt > max_font:
                        max_font = r.font_size_pt

        for s in shapes:
            if s.text:
                all_text += s.text
            for p in s.paragraphs:
                for r in p.runs:
                    all_text += r.text
                    if r.font_size_pt and r.font_size_pt > max_font:
                        max_font = r.font_size_pt

        if slide_index == total_slides - 1 and CLOSING_KEYWORDS.search(all_text):
            return "closing"

        if slide_index == 0 and element_count <= 4 and max_font >= 28:
            return "title"

        return "content"
