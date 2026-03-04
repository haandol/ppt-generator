"""PPTX 임포트 서비스: PPTX 파일 → DesignSpec 변환 오케스트레이션."""

from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation

from ppt_generator.interfaces.schemas import DesignSpec, PptxSlideSpec
from ppt_generator.tools.pptx_import.slide_reader import SlideReader

if TYPE_CHECKING:
    pass

logger = logging.getLogger(__name__)


class ImportService:
    """PPTX 파일을 DesignSpec으로 변환하는 서비스."""

    def import_from_file(self, file_path: str | Path) -> tuple[DesignSpec, list[str]]:
        """PPTX 파일 경로 → DesignSpec 변환.

        Returns:
            (DesignSpec, warnings): 변환된 디자인 스펙과 경고 목록.
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX 파일이 존재하지 않습니다: {path}")
        if not path.suffix.lower() == ".pptx":
            raise ValueError(f"PPTX 파일만 지원합니다 (.pptx): {path.suffix}")

        prs = Presentation(str(path))
        return self._convert(prs)

    def import_from_bytes(self, pptx_bytes: bytes) -> tuple[DesignSpec, list[str]]:
        """PPTX 바이트 데이터 → DesignSpec 변환.

        Returns:
            (DesignSpec, warnings): 변환된 디자인 스펙과 경고 목록.
        """
        prs = Presentation(BytesIO(pptx_bytes))
        return self._convert(prs)

    def _convert(self, prs: Presentation) -> tuple[DesignSpec, list[str]]:
        """Presentation 객체 → DesignSpec 변환."""
        scale_x, scale_y = SlideReader.compute_scale(prs)
        reader = SlideReader(scale_x=scale_x, scale_y=scale_y)

        total_slides = len(prs.slides)
        if total_slides == 0:
            raise ValueError("PPTX 파일에 슬라이드가 없습니다.")

        slides: list[PptxSlideSpec] = []
        warnings: list[str] = []

        for idx, slide in enumerate(prs.slides):
            spec = reader.read_slide(slide, idx, total_slides)
            slides.append(spec)
            logger.info(
                "슬라이드 %d/%d 추출 완료 (type=%s, textboxes=%d, shapes=%d, images=%d)",
                idx + 1,
                total_slides,
                spec.slide_type,
                len(spec.textboxes),
                len(spec.shapes),
                len(spec.images),
            )

        design_spec = DesignSpec(slides=slides)
        logger.info("PPTX 임포트 완료: %d 슬라이드", total_slides)
        return design_spec, warnings
