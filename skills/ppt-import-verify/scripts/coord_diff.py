#!/usr/bin/env python3
"""원본 PPTX 의 유효 좌표와 임포트된 design_spec 을 수치 비교한다.

사용:
    uv run python skills/ppt-import-verify/scripts/coord_diff.py "<원본.pptx>" <project_id> \
        [--slides 1,2,3] [--pos-tol 4] [--font-tol 1]

원본 slide 를 프로젝트의 SlideReader 로 파싱한 결과(placeholder 상속·그룹 좌표 변환·
차트 변환 포함)를 "기대값"으로, 저장된 design_spec/slide_NN.json 을 "실제값"으로 보고
슬라이드별 요소 개수와 텍스트박스 좌표/폰트/autofit/줄간격 차이를 표로 출력한다.

SlideReader 를 기대값 기준으로 쓰는 이유: 이미 검증된 결정론적 추출기이므로, 저장된
스펙이 최신 코드와 어긋나면(예: 서버가 옛 코드로 임포트) 그 드리프트를 드러낸다.
원본 렌더 PNG 와의 시각 비교(3단계)와 함께 보면 근본 원인 위치를 좁힐 수 있다.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path


def _project_dir(pid: str) -> Path:
    p = Path(pid)
    if p.is_absolute() and p.exists():
        return p
    return Path.home() / ".ppt-generator" / pid


def _first_text(paras) -> str:
    for para in paras:
        runs = para.get("runs") if isinstance(para, dict) else getattr(para, "runs", [])
        for r in runs or []:
            t = r.get("text") if isinstance(r, dict) else getattr(r, "text", "")
            if t and t.strip():
                return t.strip()[:24]
    return ""


def _tb_row(tb, is_spec: bool) -> dict:
    """텍스트박스를 비교용 dict 로 정규화 (spec=json dict, expected=dataclass)."""
    g = (lambda k: tb.get(k)) if is_spec else (lambda k: getattr(tb, k, None))
    paras = g("paragraphs") or []
    first_run_pt = None
    for para in paras:
        runs = para.get("runs") if is_spec else getattr(para, "runs", [])
        for r in runs:
            first_run_pt = (
                r.get("font_size_pt") if is_spec else getattr(r, "font_size_pt", None)
            )
            break
        break
    return {
        "text": _first_text(paras),
        "left": round(g("left_px") or 0, 1),
        "top": round(g("top_px") or 0, 1),
        "width": round(g("width_px") or 0, 1),
        "height": round(g("height_px") or 0, 1),
        "font_pt": first_run_pt,
        "autofit": g("autofit"),
        "linespacing": g("line_spacing_pt"),
    }


def _match(exp_rows, act_rows):
    """텍스트 첫 24자로 기대/실제 텍스트박스를 매칭. (matched pairs, only_exp, only_act)."""
    act_by_text = {}
    for r in act_rows:
        act_by_text.setdefault(r["text"], []).append(r)
    pairs = []
    only_exp = []
    for e in exp_rows:
        bucket = act_by_text.get(e["text"])
        if bucket:
            pairs.append((e, bucket.pop(0)))
        else:
            only_exp.append(e)
    only_act = [r for rows in act_by_text.values() for r in rows]
    return pairs, only_exp, only_act


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("project_id")
    ap.add_argument("--slides", default="")
    ap.add_argument("--pos-tol", type=float, default=4.0)
    ap.add_argument("--font-tol", type=float, default=1.0)
    args = ap.parse_args()

    from pptx import Presentation

    from ppt_generator.tools.pptx_import.slide_reader import SlideReader

    project_dir = _project_dir(args.project_id)
    spec_dir = project_dir / "design_spec"

    prs = Presentation(args.pptx)
    total = len(prs.slides)
    reader = SlideReader(1.0, 1.0, prs)

    slides = (
        [int(s) for s in args.slides.split(",") if s.strip()]
        if args.slides
        else list(range(1, total + 1))
    )

    grand = {"tb_pairs": 0, "tb_ok": 0, "missing": 0, "extra": 0}

    for n in slides:
        idx = n - 1
        expected = reader.read_slide(prs.slides[idx], idx, total)
        spec_path = spec_dir / f"slide_{n:02d}.json"
        if not spec_path.exists():
            print(f"\n=== slide {n}: design_spec 없음 ({spec_path.name}) ===")
            continue
        spec = json.loads(spec_path.read_text())

        exp_rows = [_tb_row(tb, is_spec=False) for tb in expected.textboxes]
        act_rows = [_tb_row(tb, is_spec=True) for tb in spec.get("textboxes", [])]

        exp_shapes = len(expected.shapes)
        act_shapes = len(spec.get("shapes", []))
        exp_imgs = len(expected.images)
        act_imgs = len(spec.get("images", []))

        print(f"\n=== slide {n} ===")
        print(
            f"  요소 개수  textbox 기대{len(exp_rows)}/실제{len(act_rows)}  "
            f"shape 기대{exp_shapes}/실제{act_shapes}  image 기대{exp_imgs}/실제{act_imgs}"
        )
        if exp_shapes != act_shapes:
            print("  [shape] 개수 불일치 → chart/stroke/missing 카테고리 의심")

        pairs, only_exp, only_act = _match(exp_rows, act_rows)
        for e, a in pairs:
            grand["tb_pairs"] += 1
            diffs = []
            for key in ("left", "top", "width", "height"):
                d = abs((e[key] or 0) - (a[key] or 0))
                if d > args.pos_tol:
                    diffs.append(f"{key}Δ{d:.0f}px")
            if (
                e["font_pt"]
                and a["font_pt"]
                and abs(e["font_pt"] - a["font_pt"]) > args.font_tol
            ):
                diffs.append(f"fontΔ{abs(e['font_pt'] - a['font_pt']):.0f}pt(→inherit)")
            if e["autofit"] != a["autofit"]:
                diffs.append(f"autofit {a['autofit']}→{e['autofit']}")
            # 줄간격: 기대 있음/실제 None → linespacing 카테고리
            if e["linespacing"] and not a["linespacing"]:
                diffs.append("linespacing 누락(→linespacing)")
            if diffs:
                print(f"  ~ {e['text']!r}: " + ", ".join(diffs))
            else:
                grand["tb_ok"] += 1
        for e in only_exp:
            grand["missing"] += 1
            print(f"  - 누락 textbox: {e['text']!r} (→missing/inherit)")
        for a in only_act:
            grand["extra"] += 1
            print(f"  + 초과 textbox: {a['text']!r}")

    print("\n=== 요약 ===")
    if grand["tb_pairs"]:
        rate = 100 * grand["tb_ok"] / grand["tb_pairs"]
        print(
            f"  매칭 textbox {grand['tb_pairs']}개 중 좌표/폰트 일치 {grand['tb_ok']}개 "
            f"({rate:.0f}%)"
        )
    print(f"  누락 {grand['missing']}개, 초과 {grand['extra']}개")
    print("  ※ shape 개수 불일치·아이콘/차트 누락은 시각 비교(3단계)와 함께 판단할 것.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
